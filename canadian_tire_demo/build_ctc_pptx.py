"""Build CTC Copilot in Power BI demo deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUT = Path(__file__).parent / "CTC_Merch_Copilot_Demo.pptx"

# CTC red + neutral palette
CTC_RED  = RGBColor(0xCA, 0x1A, 0x22)
NAVY     = RGBColor(0x1B, 0x2C, 0x5C)
GRAY     = RGBColor(0x5A, 0x6B, 0x82)
LIGHTBG  = RGBColor(0xF5, 0xF6, 0xF8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x10, 0x77, 0x52)
AMBER    = RGBColor(0xCC, 0x8A, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w:
            s.line.width = line_w
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = anchor
    tb.text_frame.margin_left = Emu(0)
    tb.text_frame.margin_right = Emu(0)
    tb.text_frame.margin_top = Emu(0)
    tb.text_frame.margin_bottom = Emu(0)
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=NAVY,
                bullet="•", line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        if isinstance(it, tuple):
            head, sub = it
            r = p.add_run()
            r.text = f"{bullet}  {head}"
            r.font.name = "Segoe UI"
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = color
            r2 = p.add_run()
            r2.text = f"  — {sub}"
            r2.font.name = "Segoe UI"
            r2.font.size = Pt(size - 1)
            r2.font.bold = False
            r2.font.color.rgb = GRAY
        else:
            r = p.add_run()
            r.text = f"{bullet}  {it}"
            r.font.name = "Segoe UI"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None, accent=CTC_RED):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill=WHITE)
    add_rect(slide, 0, Inches(0.85), SLIDE_W, Emu(38100), fill=accent)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
             title, size=24, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.35),
                 subtitle, size=11, color=GRAY)
    add_text(slide, Inches(11.5), Inches(0.25), Inches(1.7), Inches(0.35),
             "Microsoft Canada", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ----- SLIDE 1: Title --------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, Inches(0.6), Inches(3.3), Inches(1.0), Inches(0.10), fill=CTC_RED)
add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(0.6),
         "Copilot in Power BI", size=20, color=WHITE)
add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(1.0),
         "Merch Performance Demo for Canadian Tire", size=40, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.5),
         "Natural-language insights for buyers and merchants",
         size=18, color=RGBColor(0xCC, 0xD3, 0xDE))
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "Microsoft Canada | Friday May 29, 2026 | M365CopilotEnablement",
         size=11, color=RGBColor(0x9A, 0xA8, 0xBB))

# ----- SLIDE 2: Agenda -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Agenda", "What we will cover today")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5),
            [
              ("Why Copilot in Power BI",         "How natural-language consumption fits CTC's Merch workflows"),
              ("How it works",                    "Direct Lake + semantic model with rich descriptions for accurate answers"),
              ("Live demo (Merch synthetic data)", "Sales, in-season demand, connected inventory across categories and SKUs"),
              ("Customer questions, answered",    "Walking through prompts CTC prepared, live, in front of you"),
              ("Governance and roadmap",          "PII handling, data residency, conversation retention, Purview DLP"),
              ("Next steps",                      "Path from this demo to enterprise rollout at Canadian Tire"),
            ], size=16)

# ----- SLIDE 3: Audience and outcome ----------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Who this is for, what it does", "Built for the Merch buyer and analyst persona")
# 2 columns
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), fill=LIGHTBG)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4), "Audience", size=16, bold=True, color=CTC_RED)
add_bullets(s, Inches(0.7), Inches(2.05), Inches(5.6), Inches(4.5),
            [
              ("Merch buyers and analysts",  "Make day-to-day assortment, replenishment, markdown decisions"),
              ("Category and fineline leads", "Track YoY, EGM, in-season demand vs supply"),
              ("PBI Product & Platform team", "Evaluate enterprise rollout patterns and governance posture"),
            ], size=13)

add_rect(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(5.6), fill=LIGHTBG)
add_text(s, Inches(7.05), Inches(1.55), Inches(5.6), Inches(0.4), "Outcomes", size=16, bold=True, color=CTC_RED)
add_bullets(s, Inches(7.05), Inches(2.05), Inches(5.6), Inches(4.5),
            [
              ("Faster insight",            "Ask in English, no Power Query, no DAX"),
              ("Confident decisions",       "Answers grounded on the trusted semantic model"),
              ("Consistent definitions",    "POS, RVS, EGM, WoS, Lost Sales follow business glossary"),
              ("Scalable adoption",         "Same model serves buyer Q&A, dashboards, and Data Agent surfaces"),
            ], size=13)

# ----- SLIDE 4: Architecture -------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "How it works", "Direct Lake on OneLake + semantic model + Copilot")

# layered architecture
def lane(y, color, title, items):
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), fill=color)
    add_text(s, Inches(0.7), y + Emu(60000), Inches(3.0), Inches(0.4),
             title, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.8), y + Emu(60000), Inches(8.8), Inches(0.4),
             items, size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


lane(Inches(1.5), NAVY,    "Consumption",        "Copilot in Power BI  |  Smart Narrative  |  Q&A  |  Fabric Data Agent (Teams + Portal)")
lane(Inches(2.65), GRAY,   "Semantic model",     "ctc_merch (Direct Lake on OneLake) - measure descriptions, synonyms, table semantics for the Data Agent")
lane(Inches(3.80), CTC_RED, "OneLake (Delta)",    "dim_sku, dim_vendor, dim_season, dim_date, fact_sku_performance, fact_in_season, fact_connected_inventory")
lane(Inches(4.95), RGBColor(0x4F,0x55,0x60), "Source data", "Customer-prepared synthetic Merch datasets + blended dimensions (Faker)")
add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
         "No data movement at query time. Copilot resolves natural language directly against the model.",
         size=12, color=GRAY)

# ----- SLIDE 5: Data + model -------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "The data behind the demo", "Customer-prepared Merch datasets, modeled for natural language")
# left: data tiles
def tile(x, y, w, h, t, sub):
    add_rect(s, x, y, w, h, fill=LIGHTBG)
    add_text(s, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
             t, size=13, bold=True, color=CTC_RED)
    add_text(s, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), h - Inches(0.7),
             sub, size=11, color=NAVY)


tile(Inches(0.5),  Inches(1.4), Inches(4.0), Inches(2.0),
     "SKU Performance",
     "POS (Point of Sale) $ + units, RVS (Retail Value of Shipments), EGM (Enterprise Gross Margin) $ and %, QTD / YTD / R12 timeframes - 68 SKUs across 9 categories")
tile(Inches(4.65), Inches(1.4), Inches(4.0), Inches(2.0),
     "In-Season Shipment",
     "YTD POS vs YTD Ship, Retail and Corp inventory, POS and Ship forecast, Open PO, Planned Purchase, Weeks of Supply")
tile(Inches(8.80), Inches(1.4), Inches(4.0), Inches(2.0),
     "Connected Inventory",
     "Corp / Retail inventory TY vs LY, R8 sales velocity, Lost Sales %, Vendor Fill Rate %, DDF units")

# bottom: enabling features
add_rect(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.3), fill=LIGHTBG)
add_text(s, Inches(0.7), Inches(3.85), Inches(12), Inches(0.4),
         "What we added for Copilot quality", size=14, bold=True, color=CTC_RED)
add_bullets(s, Inches(0.7), Inches(4.3), Inches(12), Inches(2.7),
            [
              "Rich `description` annotation on every measure and key column (POS, RVS, EGM, WoS, Lost Sales, Fill Rate, R8, R12)",
              "Synonyms baked in (supplier vs vendor, item vs SKU, stockout rate vs lost sales)",
              "Display folders group measures by Sales, Shipments, Profitability, Inventory, WoS, Execution",
              "Customer's question set seeded into Copilot examplePrompts.json - shows up as starter prompts in the chat pane",
            ], size=12)

# ----- SLIDE 6: Live demo flow ----------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Live demo flow", "Six prompts, six minutes, business outcomes you can act on")
flow = [
    ("Summarize this page",                    "Press the Copilot button. Get an executive narrative of POS / RVS / EGM for the page in context."),
    ("Category snapshot",                       "\"Summary of POS, RVS, EGM for Kitchen & Small Appliances across QTD, YTD, R12\""),
    ("Top SKUs by EGM",                         "\"Top 10 SKUs by EGM dollars with YoY change and margin %\""),
    ("Fineline YoY comparison",                 "\"Compare YoY performance for Air Fryers vs Cookware Sets\" - includes POS, RVS, EGM %"),
    ("Overstock vs understock signal",          "\"SKUs with weeks of supply > 18 but lost sales < 2%\" - markdown candidates"),
    ("Supply-constrained growth",               "\"SKUs with lost sales > 5% AND fill rate < 85%\" - supply risk on growing SKUs"),
]
y = Inches(1.45)
for i, (head, sub) in enumerate(flow, 1):
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.85), fill=LIGHTBG)
    add_text(s, Inches(0.7), y + Inches(0.1), Inches(0.4), Inches(0.6),
             f"{i:02d}", size=18, bold=True, color=CTC_RED, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.3), y + Inches(0.05), Inches(3.3), Inches(0.4),
             head, size=13, bold=True, color=NAVY)
    add_text(s, Inches(1.3), y + Inches(0.42), Inches(11.3), Inches(0.4),
             sub, size=11, color=GRAY)
    y += Inches(0.92)

# ----- SLIDE 7: CTC Merch Data Agent ----------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "CTC Merch Data Agent", "Conversational analytics outside the report - in Teams, in the portal")

# Left side: what it is + where it lives
add_rect(s, Inches(0.5), Inches(1.4), Inches(5.5), Inches(5.6), fill=LIGHTBG)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.1), Inches(0.4),
         "Same model, different surface", size=15, bold=True, color=CTC_RED)
add_bullets(s, Inches(0.7), Inches(2.05), Inches(5.1), Inches(2.6),
            [
              ("Grounded on ctc_merch",      "Same Direct Lake model that powers the report"),
              ("Merch vocabulary built in",  "POS, EGM, RVS, WoS, fill rate, lost sales, fineline"),
              ("Use cases pre-seeded",       "Top SKUs by EGM, at-risk SKUs, markdown candidates, vendor summary"),
              ("Lives where buyers work",    "Fabric portal today; Teams + Copilot Studio next"),
            ], size=12)
add_text(s, Inches(0.7), Inches(5.4), Inches(5.1), Inches(0.4),
         "Status: published, draft + published stages bound, 5 relationships indexed",
         size=10, color=GRAY)
add_text(s, Inches(0.7), Inches(5.85), Inches(5.1), Inches(0.4),
         "7 tables - 75 columns - 64 measures - all annotated",
         size=10, color=GRAY)

# Right side: 10 starter prompts
add_rect(s, Inches(6.15), Inches(1.4), Inches(6.7), Inches(5.6), fill=LIGHTBG)
add_text(s, Inches(6.35), Inches(1.55), Inches(6.3), Inches(0.4),
         "10 Starter Questions buyers can run today", size=15, bold=True, color=CTC_RED)
prompts = [
    "One-paragraph merch briefing - POS YoY, EGM %, growth and drag",
    "Top 10 SKUs by EGM dollars with POS YoY and EGM %",
    "Air Fryers vs Cookware Sets - POS YoY and EGM %",
    "SKUs with lost sales > 5% AND fill rate < 85%",
    "WoS > 18 AND lost sales < 2% - markdown candidates",
    "Canvas Outdoor vendor - POS, fill rate, EGM contribution",
    "New SKUs (New_SKU_Flag = Yes) driving the most POS growth",
    "Worst POS YoY % decline by fineline and its fill rate",
    "Vendors with worst fill rate but highest planned purchase",
    "Demand vs supply gap by category - where to expedite",
]
y = Inches(2.05)
for i, q in enumerate(prompts, 1):
    add_text(s, Inches(6.35), y, Inches(0.4), Inches(0.3),
             f"{i:02d}", size=10, bold=True, color=CTC_RED)
    add_text(s, Inches(6.75), y, Inches(5.9), Inches(0.3),
             q, size=10, color=NAVY)
    y += Inches(0.46)

# ----- SLIDE 8: Headline findings -------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Headline findings from the demo data", "What Copilot surfaces in seconds")


def kpi_card(x, y, w, h, value, label, color=CTC_RED):
    add_rect(s, x, y, w, h, fill=LIGHTBG)
    add_text(s, x, y + Inches(0.25), w, Inches(0.7), value,
             size=24, bold=True, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y + Inches(1.05), w, Inches(0.4), label,
             size=11, color=GRAY, align=PP_ALIGN.CENTER)


kpi_card(Inches(0.5),  Inches(1.4), Inches(3.0), Inches(1.6), "+27.9%", "POS YoY (Total)")
kpi_card(Inches(3.65), Inches(1.4), Inches(3.0), Inches(1.6), "42.3%",  "EGM % TY", color=GREEN)
kpi_card(Inches(6.80), Inches(1.4), Inches(3.0), Inches(1.6), "4.06%",  "Avg Lost Sales", color=AMBER)
kpi_card(Inches(9.95), Inches(1.4), Inches(3.0), Inches(1.6), "89.0%",  "Vendor Fill Rate")

add_text(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.4),
         "Insights merchants can act on", size=14, bold=True, color=CTC_RED)
add_bullets(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(3.5),
            [
              ("Air Fryers +88.7% YoY POS",            "Strongest growth fineline - protect supply, expand assortment"),
              ("NK Dual Zone Air Fryer 8L losing 13.7% of demand",  "Vendor fill rate only 72.7% - immediate supply intervention"),
              ("9 SKUs WoS > 18 with low lost sales",  "Classic markdown / promo candidates - sell-through stalled"),
              ("OL-BBQ-003 (Woodfire Grill) +1671% YoY", "New SKU with strong demand but fill rate 71% - over-rotate on PO"),
              ("Canvas Outdoor vendor +2.9% YoY, $7.4M, 39.9% EGM, 89.1% fill",  "Anchor vendor performing steadily"),
            ], size=12)

# ----- SLIDE 8: Governance --------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Governance and trust", "Built for enterprise rollout")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
            [
              ("PII handling",            "Copilot in PBI does not send PII outside the M365 service boundary. Sensitivity labels and Purview classification carry through."),
              ("Data residency",          "Canadian Tire is on Canada-resident Fabric capacity. Copilot processing stays within the M365 geo of your tenant."),
              ("Conversation retention",  "Prompt and response history governed by M365 admin retention policies. No model training on your data."),
              ("Purview DLP",             "Existing DLP policies apply to Copilot interactions, including blocked sharing of sensitive content."),
              ("Responsible AI",          "Aligned with Microsoft Responsible AI Standard - groundedness, safety filters, attribution."),
              ("Audit and observability", "Copilot usage logged in M365 audit log, viewable via Purview eDiscovery."),
            ], size=13)

# ----- SLIDE 9: Roadmap -----------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Roadmap", "Where Copilot in Power BI is going")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
            [
              ("GA today",                 "Summarize page, Q&A, Smart Narrative, Find Insights, anomaly explanations"),
              ("Recent additions",         "Cross-report grounding, custom semantic model instructions, multi-turn Q&A"),
              ("Near-term (next 90 days)", "Agent-style follow-ups, deeper drill into measure lineage, custom example prompts at workspace scope"),
              ("Integration story",        "Fabric Data Agent, Copilot Studio agents, Teams + Outlook + Word surfaces"),
              ("CTC-specific opportunities", "Workspace-scoped Copilot for Merch, finance-team scope (separate effort), supply-chain agent on Connected Inventory"),
            ], size=13)

# ----- SLIDE 10: Where to next ----------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Recommended next steps", "From demo to enterprise rollout")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
            [
              ("This week",     "Enable Copilot in Power BI for the Merch team pilot workspace (IT ticket - CTC IT noted ongoing enablement)"),
              ("Within 2 weeks", "Apply the descriptions / synonyms / example prompts pattern to the migrated Merch semantic models"),
              ("Within 4 weeks", "Pilot with 10-15 buyers, capture prompt patterns, refine examplePrompts.json"),
              ("ARB-2",          "Use this pilot evidence + governance posture for ARB-2 approval"),
              ("Broader rollout", "Cascade to Finance (separate effort), Supply Chain, Stores once Merch validates"),
              ("Microsoft Canada support", "Engineering hands-on for the first 30 days; weekly office hours for buyer Q&A patterns"),
            ], size=13)

# ----- SLIDE 11: Closing ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, Inches(0.6), Inches(3.0), Inches(1.0), Inches(0.10), fill=CTC_RED)
add_text(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.0),
         "Questions?", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.6),
         "Let's talk about how Copilot fits Merch at Canadian Tire",
         size=18, color=RGBColor(0xCC, 0xD3, 0xDE))
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "anuragdh@microsoft.com  |  Microsoft Canada Data Specialist",
         size=11, color=RGBColor(0x9A, 0xA8, 0xBB))

prs.save(OUT)
print(f"Wrote {OUT}")
