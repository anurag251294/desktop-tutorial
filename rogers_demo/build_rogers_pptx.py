"""Build Rogers Network Anomaly Detection on Fabric demo deck.

12 slides covering the same arc as the CTC deck but for network operations:
  1. Title
  2. The problem - network anomalies cost minutes and customers
  3. Architecture lanes
  4. Data model schema
  5. Four anomaly stories seeded in the data
  6. Page 1 - Network Health Overview
  7. Page 2 - Anomaly Deep-Dive
  8. Page 3 - Customer Impact
  9. Rogers Network Data Agent (same model, conversational)
 10. Live demo path
 11. Next steps
 12. Closer
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUT = Path(__file__).parent / "Rogers_Network_Copilot_Demo.pptx"

ROGERS_RED = RGBColor(0xDA, 0x29, 0x1C)
NAVY    = RGBColor(0x1B, 0x2C, 0x5C)
GRAY    = RGBColor(0x5A, 0x6B, 0x82)
LIGHTBG = RGBColor(0xF5, 0xF6, 0xF8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x10, 0x77, 0x52)
AMBER   = RGBColor(0xCC, 0x8A, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, Emu(0))
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=NAVY,
                bullet="-", line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, Emu(0))
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        if isinstance(it, tuple):
            head, sub = it
            r = p.add_run(); r.text = f"{bullet}  {head}"
            r.font.name = "Segoe UI"; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
            r2 = p.add_run(); r2.text = f"  - {sub}"
            r2.font.name = "Segoe UI"; r2.font.size = Pt(size-1); r2.font.color.rgb = GRAY
        else:
            r = p.add_run(); r.text = f"{bullet}  {it}"
            r.font.name = "Segoe UI"; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.55), fill=ROGERS_RED)
    add_text(slide, Inches(0.5), Inches(0.08), Inches(12), Inches(0.42),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.35),
                 subtitle, size=14, color=GRAY)


def add_footer(slide, text="Rogers Network Anomaly Detection on Microsoft Fabric"):
    add_text(slide, Inches(0.5), Inches(7.15), Inches(12), Inches(0.3),
             text, size=10, color=GRAY)


# ---- SLIDE 1: TITLE -------------------------------------------------------

s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, Inches(0), Inches(3.4), SLIDE_W, Inches(0.05), fill=ROGERS_RED)
add_text(s, Inches(0.6), Inches(2.2), Inches(12.0), Inches(0.8),
         "Rogers", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.0), Inches(12.0), Inches(0.6),
         "Network Anomaly Detection on Microsoft Fabric",
         size=26, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.7), Inches(12.0), Inches(0.5),
         "Direct Lake + Copilot + Data Agent | Executive briefing demo",
         size=16, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.4),
         "Microsoft Canada  |  Rogers Network",
         size=12, color=RGBColor(0xCF, 0xD8, 0xE3))


# ---- SLIDE 2: WHY -------------------------------------------------------

s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Why network anomaly detection",
              "Minutes-to-detect drives MTTR; MTTR drives customer impact; customer impact drives churn")

add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4),
         "The pattern we see:", size=14, bold=True, color=NAVY)

cards = [
    ("Network signal lives in silos",
     "RAN OSS, transport NMS, core EMS, alarms, traffic, customer care - five systems, five lenses, five tools."),
    ("Anomalies are noticed late",
     "A sector outage at 16:00 becomes a ticket spike at 09:00 the next day, and a churn risk by Friday."),
    ("Asking 'why' takes 5 tools",
     "By the time the NOC reconstructs the story across systems, the next anomaly is already happening."),
]
y = 1.6
for title, desc in cards:
    add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.4), fill=LIGHTBG)
    add_rect(s, Inches(0.5), Inches(y), Inches(0.1), Inches(1.4), fill=ROGERS_RED)
    add_text(s, Inches(0.85), Inches(y+0.15), Inches(11.5), Inches(0.4),
             title, size=15, bold=True, color=NAVY)
    add_text(s, Inches(0.85), Inches(y+0.55), Inches(11.5), Inches(0.8),
             desc, size=12, color=GRAY)
    y += 1.6

add_footer(s)


# ---- SLIDE 3: ARCHITECTURE ------------------------------------------------

s = prs.slides.add_slide(BLANK)
add_title_bar(s, "One platform, three lanes, one semantic model",
              "Microsoft Fabric replaces five separate stacks with one OneLake-backed brain")

lane_y = 1.4
lane_h = 5.3
lane_w_in = 4.2
for i, (label, color, items) in enumerate([
    ("Sources",
     RGBColor(0x37, 0x4F, 0xA0),
     ["RAN OSS (Ericsson / Nokia / Samsung)",
      "Transport NMS",
      "Core EMS (PCRF / IMS)",
      "Power / DC telemetry",
      "Care contact-centre data",
      "Subscriber & segment lookups"]),
    ("OneLake + Direct Lake semantic model",
     ROGERS_RED,
     ["Lakehouse holds 4 facts + 6 dims",
      "Direct Lake mode - no copies",
      "TMDL with measure descriptions",
      "Anomaly measures: cells <99%, PRB>85%, latency >60ms, MTTR",
      "Network Health Score (composite)"]),
    ("Consumption surfaces",
     RGBColor(0x10, 0x77, 0x52),
     ["Power BI exec report",
      "  - 3 pages: Overview, Deep-Dive, Customer Impact",
      "Fabric Data Agent (Teams + M365 Copilot)",
      "Copilot in Power BI (per-visual narrative)",
      "Real-time KPI cards"]),
]):
    x = Inches(0.4 + i * (lane_w_in + 0.1))
    add_rect(s, x, Inches(lane_y), Inches(lane_w_in), Inches(0.5), fill=color)
    add_text(s, x, Inches(lane_y), Inches(lane_w_in), Inches(0.5),
             label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, Inches(lane_y + 0.5), Inches(lane_w_in), Inches(lane_h - 0.5), fill=LIGHTBG)
    add_bullets(s, x + Emu(180000), Inches(lane_y + 0.65),
                Inches(lane_w_in - 0.3), Inches(lane_h - 0.7),
                items, size=12)

add_footer(s)


# ---- SLIDE 4: DATA MODEL --------------------------------------------------

s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Data model - rogers_net semantic model",
              "Star schema in Direct Lake: 4 facts share the cell_id key")

# Fact tables (centre)
facts = [
    ("fact_cell_kpi", "Hourly per-cell radio KPIs", 4.7),
    ("fact_alarms", "Discrete alarm events + MTTR", 5.6),
    ("fact_traffic", "Daily voice/data/SMS rollups", 1.9),
    ("fact_customer_impact", "Daily tickets x segment", 2.8),
]
for name, desc, y in facts:
    add_rect(s, Inches(5.0), Inches(y), Inches(3.3), Inches(0.8), fill=ROGERS_RED)
    add_text(s, Inches(5.0), Inches(y + 0.05), Inches(3.3), Inches(0.35),
             name, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.0), Inches(y + 0.42), Inches(3.3), Inches(0.32),
             desc, size=10, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Dims (left + right)
dims_left = [
    ("dim_date",   "calendar"),
    ("dim_hour",   "hour-of-day"),
    ("dim_site",   "tower master"),
]
dims_right = [
    ("dim_cell",              "sector/cell master"),
    ("dim_alarm_type",        "alarm catalog"),
    ("dim_customer_segment",  "segments"),
]
for items, x_in in [(dims_left, 0.5), (dims_right, 9.3)]:
    for i, (name, desc) in enumerate(items):
        y = 2.0 + i * 1.6
        add_rect(s, Inches(x_in), Inches(y), Inches(3.2), Inches(0.8), fill=NAVY)
        add_text(s, Inches(x_in), Inches(y + 0.05), Inches(3.2), Inches(0.35),
                 name, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(x_in), Inches(y + 0.42), Inches(3.2), Inches(0.32),
                 desc, size=10, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
         "12 relationships  |  All many-to-one  |  Direct Lake mode (no model refresh required)",
         size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ---- SLIDE 5: FOUR ANOMALIES ----------------------------------------------

s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Four anomaly stories seeded in the data",
              "Each one becomes a Copilot question; each answer surfaces the same model")

stories = [
    ("A. Site outage (S-TOR-009)",
     "3 cells go hard down 16:00-19:00 on 2026-05-25. Availability drops to ~15%. Customer tickets spike +470% the following morning across Consumer-Postpaid + Enterprise."),
    ("B. Vendor latency creep",
     "Nokia cells across BC show RTT drifting from 22ms to 55ms starting 2026-05-22 - no critical alarms, classic systemic vendor signal."),
    ("C. Stadium-event congestion",
     "Downtown Toronto cells hit 96% PRB utilization 18:00-22:00 on 2026-05-24. Throughput per user drops 60%. Traffic surge in dim_hour x dim_site."),
    ("D. Chronic cell degradation",
     "C-VAN-N-002 availability declines linearly 99% -> 89% over the period - no alarms - hardware-fault candidate."),
]
y = 1.3
for i, (title, desc) in enumerate(stories):
    add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.35), fill=LIGHTBG)
    add_rect(s, Inches(0.5), Inches(y), Inches(0.1), Inches(1.35), fill=ROGERS_RED)
    add_text(s, Inches(0.85), Inches(y+0.12), Inches(11.5), Inches(0.4),
             title, size=15, bold=True, color=NAVY)
    add_text(s, Inches(0.85), Inches(y+0.5), Inches(11.5), Inches(0.8),
             desc, size=12, color=GRAY)
    y += 1.5

add_footer(s)


# ---- SLIDES 6-8: REPORT PAGE MOCKUPS --------------------------------------

def report_mock_slide(title, subtitle, kpis, narrative_lines, prompts):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, title, subtitle)

    # 4 KPI cards
    for i, (label, value, sub) in enumerate(kpis):
        x = Inches(0.4 + i * 3.2)
        add_rect(s, x, Inches(1.25), Inches(3.0), Inches(1.2), fill=WHITE, line=RGBColor(0xE5,0xE7,0xEB), line_w=Pt(0.75))
        add_text(s, x, Inches(1.32), Inches(3.0), Inches(0.35), label, size=10, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(1.62), Inches(3.0), Inches(0.55), value, size=22, bold=True, color=ROGERS_RED, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.10), Inches(3.0), Inches(0.3), sub, size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # narrative box
    add_rect(s, Inches(0.4), Inches(2.65), Inches(6.0), Inches(3.6), fill=LIGHTBG)
    add_text(s, Inches(0.55), Inches(2.75), Inches(5.7), Inches(0.4),
             "AI Narrative", size=12, bold=True, color=NAVY)
    add_bullets(s, Inches(0.55), Inches(3.15), Inches(5.7), Inches(3.0), narrative_lines, size=11)

    # Copilot prompt box
    add_rect(s, Inches(6.6), Inches(2.65), Inches(6.3), Inches(3.6), fill=RGBColor(0xFF, 0xF6, 0xF6), line=ROGERS_RED, line_w=Pt(0.75))
    add_text(s, Inches(6.75), Inches(2.75), Inches(6.0), Inches(0.35),
             "Try these Copilot prompts", size=12, bold=True, color=ROGERS_RED)
    add_bullets(s, Inches(6.75), Inches(3.15), Inches(6.0), Inches(3.0), prompts, size=11)

    # bottom strip - slicers
    for i, lab in enumerate(["Province", "Vendor", "Technology"]):
        add_rect(s, Inches(0.4 + i * 4.2), Inches(6.4), Inches(4.05), Inches(0.4), fill=WHITE, line=RGBColor(0xE5,0xE7,0xEB), line_w=Pt(0.5))
        add_text(s, Inches(0.4 + i * 4.2), Inches(6.42), Inches(4.05), Inches(0.36),
                 lab, size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s)


report_mock_slide(
    "Page 1 - Network Health Overview",
    "Direct Lake | RAN KPIs x Alarms x Customer Impact | Copilot exec demo",
    [("Avg Availability %", "99.3%", "Last 7d"),
     ("Avg DL Throughput", "412 Mbps", "Across cells"),
     ("Cells <99% Avail", "21", "Anomaly headline"),
     ("Tickets Opened", "17.4K", "Across segments")],
    ["Availability holds at 99.3% on average but 21 cells fall below threshold",
     "Critical alarms cluster in ON (Tor sites) and BC",
     "Tickets ticked up +18% week-on-week, driven by Postpaid + Enterprise"],
    ["Brief me on last 7 days network health",
     "Sites with availability below 95%",
     "Throughput by vendor and province",
     "Top 10 most congested cells",
     "Which segment opened the most tickets?"]
)

report_mock_slide(
    "Page 2 - Anomaly Deep-Dive",
    "Availability trend x Throughput by vendor x Latency anomalies x Congestion hotspots",
    [("Worst-Hour Availability", "12.4%", "S-TOR-009 outage"),
     ("Avg Latency", "31 ms", "Nokia BC drifting up"),
     ("Avg PRB Utilization", "58%", "Spikes downtown TOR"),
     ("Avg MTTR", "42 min", "Critical-severity")],
    ["Outage hours concentrated in 16:00-19:00 window on 2026-05-25",
     "Nokia cells in BC show latency creep without firing alarms",
     "Stadium-area cells hit PRB > 85% for 4 hours on 2026-05-24 evening"],
    ["Which cells had latency above 60ms last week?",
     "Vendors with worst throughput, broken out by province?",
     "Cell-hours at >85% PRB last 7 days?",
     "Walk me through 2026-05-25 site outage",
     "Cells with steady decline - hardware-fault candidates?"]
)

report_mock_slide(
    "Page 3 - Customer Impact",
    "Tickets by segment x Sentiment x Affected customers x Repeat caller rate",
    [("Tickets Opened", "17.4K", "Period total"),
     ("Repeat Caller Rate", "24.1%", "Higher when CX bad"),
     ("Avg Sentiment", "0.54", "Postpaid lowest"),
     ("Affected Customers", "612K", "Network-driven")],
    ["Enterprise + Consumer-Postpaid open the most tickets per impacted cell",
     "Sentiment drops 0.18 in ON the day after the outage",
     "Repeat-caller rate peaks 28% in BC during latency drift"],
    ["Which segment opened the most tickets?",
     "Where did affected customers spike?",
     "Sentiment shift after the 2026-05-25 outage?",
     "Provinces with the worst repeat-caller rate?",
     "Sentiment trend last 7 days for Enterprise"]
)


# ---- SLIDE 9: DATA AGENT --------------------------------------------------

s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Rogers Network Data Agent",
              "Same model, conversational - in Teams, in M365 Copilot, in the Fabric portal")

# left panel - what it is
add_rect(s, Inches(0.5), Inches(1.2), Inches(5.6), Inches(5.5), fill=LIGHTBG)
add_text(s, Inches(0.7), Inches(1.3), Inches(5.2), Inches(0.4),
         "Same model, different surface", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(1.7), Inches(5.2), Inches(2.6),
            ["Grounded in rogers_net semantic model",
             "Inherits every measure description",
             "Telco vocab built into AI instructions (RAN, PRB, RTT, MTTR)",
             "10 starter prompts seeded for NOC + exec briefing",
             "Use in Teams chat or M365 Copilot",
             "Show data lineage for every answer"], size=11)
add_text(s, Inches(0.7), Inches(4.5), Inches(5.2), Inches(0.4),
         "Status", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(4.9), Inches(5.2), Inches(1.6),
            ["10 tables wired",
             "75+ columns + measures described",
             "12 relationships imported",
             "Published + draft stages aligned",
             "Live in Fabric portal"], size=11)

# right panel - starter prompts
add_rect(s, Inches(6.4), Inches(1.2), Inches(6.5), Inches(5.5), fill=WHITE, line=ROGERS_RED, line_w=Pt(0.75))
add_text(s, Inches(6.6), Inches(1.3), Inches(6.1), Inches(0.4),
         "10 starter prompts (NOC + exec briefing)", size=14, bold=True, color=ROGERS_RED)
prompts = [
    "1. 7-day network health briefing - biggest anomaly",
    "2. Sites with availability <95% + customer impact",
    "3. DL throughput by vendor x province",
    "4. Cells with avg latency >60ms (vendor / province)",
    "5. Top 10 most congested cells",
    "6. Walk through 2026-05-25 outage (alarms, MTTR, tix)",
    "7. Tickets per impacted cell by segment + sentiment",
    "8. Chronic availability decline candidates",
    "9. Tonight 18:00-22:00 congestion risk",
    "10. Total data volume by province x 4G/5G",
]
add_bullets(s, Inches(6.6), Inches(1.7), Inches(6.1), Inches(4.8), prompts, size=10, bullet="")

add_footer(s)


# ---- SLIDE 10: LIVE DEMO PATH ---------------------------------------------

s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Live demo path", "5 minutes, 5 questions, one story")

steps = [
    ("1. Exec snapshot",   "Open Network Health Overview - hero KPIs, narrative tells the headline"),
    ("2. Spot the anomaly", "Question to Copilot: 'Sites with availability below 95% last week?' - S-TOR-009 surfaces"),
    ("3. Drill the why",    "Page 2 - 16:00-19:00 outage on 3 cells, 3 critical alarms, MTTR 180 min"),
    ("4. Quantify customer impact","Page 3 - Postpaid + Enterprise tickets +470% the next morning"),
    ("5. Continue in Teams","Same question to Rogers Network Data Agent - same answer, in chat, with citation"),
]
y = 1.25
for h, d in steps:
    add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.0), fill=LIGHTBG)
    add_rect(s, Inches(0.5), Inches(y), Inches(0.1), Inches(1.0), fill=ROGERS_RED)
    add_text(s, Inches(0.85), Inches(y+0.1), Inches(4.5), Inches(0.4), h, size=14, bold=True, color=NAVY)
    add_text(s, Inches(5.4), Inches(y+0.15), Inches(7.4), Inches(0.75), d, size=12, color=GRAY)
    y += 1.15

add_footer(s)


# ---- SLIDE 11: NEXT STEPS -------------------------------------------------

s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Where this goes next",
              "From single-region pilot to enterprise NOC + Care + Eng integration")

# 3 columns
cols = [
    ("0-30 days",  GREEN, ["Pick one region (ON or BC) and one anomaly class",
                             "Wire real OSS / NMS extracts via Mirroring or shortcuts",
                             "Validate measure semantics with NOC SMEs",
                             "30-min daily 'briefing me' usage from VP Network"]),
    ("30-90 days",  AMBER, ["Add transport + core fact tables",
                             "Tune anomaly thresholds with real-world incident data",
                             "Embed in Teams for on-call rotations",
                             "Connect tickets + Eng tickets back to network signal"]),
    ("90+ days",  ROGERS_RED, ["Real-time (Eventhouse) for sub-minute anomaly detection",
                                 "Predictive: forecast next-day congestion + capex triggers",
                                 "Customer-journey view: churn risk from cumulative network experience",
                                 "Co-pilot in Field Service for dispatch optimization"]),
]
for i, (h, color, items) in enumerate(cols):
    x = Inches(0.4 + i * 4.3)
    add_rect(s, x, Inches(1.2), Inches(4.2), Inches(0.55), fill=color)
    add_text(s, x, Inches(1.2), Inches(4.2), Inches(0.55), h, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, Inches(1.75), Inches(4.2), Inches(5.0), fill=LIGHTBG)
    add_bullets(s, x + Emu(180000), Inches(1.9), Inches(3.9), Inches(4.7), items, size=11)

add_footer(s)


# ---- SLIDE 12: CLOSER -----------------------------------------------------

s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, Inches(0), Inches(3.5), SLIDE_W, Inches(0.05), fill=ROGERS_RED)
add_text(s, Inches(0.6), Inches(2.3), Inches(12.0), Inches(0.8),
         "Thank you", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.1), Inches(12.0), Inches(0.5),
         "Anomalies surfaced. Stories explained. Customers protected.",
         size=20, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.6), Inches(3.9), Inches(12.0), Inches(0.4),
         "One Direct Lake model | One Power BI report | One Data Agent | One platform.",
         size=14, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.4),
         "Microsoft Canada x Rogers Network", size=12, color=RGBColor(0xCF, 0xD8, 0xE3))


# ---- save -----------------------------------------------------------------

prs.save(OUT)
print(f"Saved {OUT}  ({OUT.stat().st_size:,} bytes)  {len(prs.slides)} slides")
