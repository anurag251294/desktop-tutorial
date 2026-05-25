"""Build Interac Power BI + Fabric demo deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUT = Path(__file__).parent / "Interac_PBI_Fabric_Demo.pptx"

# Interac orange + neutral
INT_ORANGE = RGBColor(0xF5, 0x82, 0x20)
NAVY       = RGBColor(0x1B, 0x2C, 0x5C)
GRAY       = RGBColor(0x5A, 0x6B, 0x82)
LIGHTBG    = RGBColor(0xF5, 0xF6, 0xF8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x10, 0x77, 0x52)
AMBER      = RGBColor(0xCC, 0x8A, 0x00)
RED        = RGBColor(0xC0, 0x3A, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = anchor
    tb.text_frame.margin_left = Emu(0)
    tb.text_frame.margin_right = Emu(0)
    tb.text_frame.margin_top = Emu(0)
    tb.text_frame.margin_bottom = Emu(0)
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=NAVY, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        if isinstance(it, tuple):
            head, sub = it
            r = p.add_run(); r.text = f"{bullet}  {head}"
            r.font.name = "Segoe UI"; r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = color
            r2 = p.add_run(); r2.text = f"  — {sub}"
            r2.font.name = "Segoe UI"; r2.font.size = Pt(size - 1)
            r2.font.color.rgb = GRAY
        else:
            r = p.add_run(); r.text = f"{bullet}  {it}"
            r.font.name = "Segoe UI"; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None, accent=INT_ORANGE):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill=WHITE)
    add_rect(slide, 0, Inches(0.85), SLIDE_W, Emu(38100), fill=accent)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
             title, size=24, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.35),
                 subtitle, size=11, color=GRAY)
    add_text(slide, Inches(11.5), Inches(0.25), Inches(1.7), Inches(0.35),
             "Microsoft Canada", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ----- SLIDE 1: Title --------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, Inches(0.6), Inches(3.3), Inches(1.0), Inches(0.10), fill=INT_ORANGE)
add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(0.6),
         "Power BI on Fabric", size=20, color=WHITE)
add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(1.0),
         "HR Analytics Demo for Interac", size=40, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.5),
         "OSFI / FINTRAC-aware workforce insights, in Direct Lake",
         size=18, color=RGBColor(0xCC, 0xD3, 0xDE))
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "Microsoft Canada | Frontier Transformation - Financial Services",
         size=11, color=RGBColor(0x9A, 0xA8, 0xBB))

# ----- SLIDE 2: Agenda -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Agenda", "What we will cover")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5),
            [
              ("Objectives and success criteria",   "Align on what \"good\" looks like for this evaluation"),
              ("Power BI on Fabric - capabilities", "Direct Lake, OneLake, governance, Copilot, Data Agent"),
              ("HR demo (synthetic Interac data)",  "Active employees, attrition, FINTRAC training, COI attestations"),
              ("RLS and persona views",             "HR Business Partner / People Manager / Compliance Officer"),
              ("Copilot and Data Agent moments",    "Natural-language questions on the live model"),
              ("Deployment, licensing, and roadmap","From demo to production at Interac"),
            ], size=16)

# ----- SLIDE 3: Why Interac, why now ----------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Why this matters for Interac", "Regulated workforce data, demanding stewardship")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
            [
              ("OSFI E-21 expectations",       "Sound operational-risk management requires timely workforce visibility for critical-function roles"),
              ("FINTRAC training compliance",  "Mandatory training completion across regulated employees - audit trail and timely attestation"),
              ("Conflict-of-interest attestations", "COI overdue >90 days is a regulator-visible risk indicator"),
              ("Talent risk in Tech and Risk",  "Regrettable attrition concentrated in mission-critical functions - early-warning needed"),
              ("Operating model",               "HR / Risk / Tech leadership want one trusted source - not 6 spreadsheets"),
            ], size=14)

# ----- SLIDE 4: Architecture -------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "How it works", "Direct Lake on OneLake + semantic model + Copilot + Data Agent")


def lane(y, color, title, items):
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), fill=color)
    add_text(s, Inches(0.7), y + Emu(60000), Inches(3.2), Inches(0.4),
             title, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.0), y + Emu(60000), Inches(8.5), Inches(0.4),
             items, size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


lane(Inches(1.5), NAVY,    "Consumption",     "Copilot in Power BI  |  HR Dashboard  |  Fabric Data Agent  |  Teams / Outlook")
lane(Inches(2.65), GRAY,   "Semantic model",  "hr_demo (Direct Lake on OneLake) - 8 KPIs, 22 measures, 3 RLS roles, Copilot example prompts")
lane(Inches(3.80), INT_ORANGE, "OneLake (Delta)", "dim_date, dim_department, dim_employee, dim_location, dim_role, 6 fact tables")
lane(Inches(4.95), RGBColor(0x4F,0x55,0x60), "Source data", "Synthetic Interac HR data - 1000 employees, OSFI / FINTRAC overlay - SEED=20260524")
add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
         "Zero data movement at query time. Copilot resolves natural language directly against the model.",
         size=12, color=GRAY)

# ----- SLIDE 5: Headline KPIs ----------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Live numbers from the demo data", "8 KPIs with traffic-light status")


def kpi(x, y, w, h, value, label, color, status_text):
    add_rect(s, x, y, w, h, fill=LIGHTBG)
    add_text(s, x, y + Inches(0.25), w, Inches(0.7), value, size=22, bold=True,
             color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y + Inches(1.05), w, Inches(0.35), label,
             size=11, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.45), w, Inches(0.3), status_text,
             size=9, color=color, align=PP_ALIGN.CENTER, bold=True)


# Row 1
kpi(Inches(0.5),  Inches(1.4), Inches(3.0), Inches(2.0), "899",   "Active Employees vs 1,000 target", AMBER, "YELLOW")
kpi(Inches(3.65), Inches(1.4), Inches(3.0), Inches(2.0), "14.5%", "Attrition Rate LTM",               AMBER, "YELLOW")
kpi(Inches(6.80), Inches(1.4), Inches(3.0), Inches(2.0), "22.5%", "Regrettable Attrition %",           AMBER, "YELLOW")
kpi(Inches(9.95), Inches(1.4), Inches(3.0), Inches(2.0), "96.3%", "FINTRAC Training Completion",       GREEN,  "GREEN")
# Row 2
kpi(Inches(0.5),  Inches(3.6), Inches(3.0), Inches(2.0), "7",     "COI Overdue 90+ Days",              RED,    "RED")
kpi(Inches(3.65), Inches(3.6), Inches(3.0), Inches(2.0), "62",    "Open Requisitions",                 AMBER, "YELLOW")
kpi(Inches(6.80), Inches(3.6), Inches(3.0), Inches(2.0), "68d",   "Avg Time to Fill",                  AMBER, "YELLOW")
kpi(Inches(9.95), Inches(3.6), Inches(3.0), Inches(2.0), "0.98",  "Comp Ratio vs Market",              AMBER, "YELLOW")

add_text(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.5),
         "Each KPI has a `kpi` TMDL block with traffic-light status against an explicit target.",
         size=12, color=GRAY)

# ----- SLIDE 6: Demo flow ---------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Live demo flow", "Persona-led walkthrough on the live model")
flow = [
    ("HR Business Partner view",   "Open HR Dashboard - 8 KPIs, traffic lights, drill from total attrition to function and department"),
    ("People Manager view (RLS)",  "Same report, View As role = People Manager - only the manager's team is visible"),
    ("Compliance Officer view",    "View As Compliance Officer - filtered to Risk & Compliance function for COI review"),
    ("Copilot - exec summary",     "Click Copilot, ask \"summarize this page\" - get narrative across all 8 KPIs"),
    ("Copilot - COI investigation","\"Which departments have COI attestations overdue 90+ days?\" - returns table + chart"),
    ("Data Agent - Teams surface", "Same questions answered inside Teams via the Fabric Data Agent - bot-style"),
]
y = Inches(1.45)
for i, (head, sub) in enumerate(flow, 1):
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.85), fill=LIGHTBG)
    add_text(s, Inches(0.7), y + Inches(0.1), Inches(0.4), Inches(0.6),
             f"{i:02d}", size=18, bold=True, color=INT_ORANGE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.3), y + Inches(0.05), Inches(3.3), Inches(0.4),
             head, size=13, bold=True, color=NAVY)
    add_text(s, Inches(1.3), y + Inches(0.42), Inches(11.3), Inches(0.4),
             sub, size=11, color=GRAY)
    y += Inches(0.92)

# ----- SLIDE 7: Copilot + Data Agent ----------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Copilot and Data Agent on the same model", "Two surfaces, one source of truth")
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), fill=LIGHTBG)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
         "Copilot in Power BI", size=16, bold=True, color=INT_ORANGE)
add_bullets(s, Inches(0.7), Inches(2.05), Inches(5.6), Inches(4.5),
            [
              ("Inside the report",          "Summarize, Q&A, Smart Narrative, Find Insights"),
              ("Author and consume",         "Same surface for the analyst and the executive"),
              ("Grounded on the model",      "No hallucinated DAX - measures from the semantic glossary"),
              ("Example prompts",            "examplePrompts.json shows starter questions to users"),
            ], size=12)

add_rect(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(5.6), fill=LIGHTBG)
add_text(s, Inches(7.05), Inches(1.55), Inches(5.6), Inches(0.4),
         "Fabric Data Agent", size=16, bold=True, color=INT_ORANGE)
add_bullets(s, Inches(7.05), Inches(2.05), Inches(5.6), Inches(4.5),
            [
              ("Outside the report",         "Conversational answers in Teams, Outlook, Web"),
              ("Same semantic model",        "Reuses descriptions and synonyms - consistent answers"),
              ("Async use cases",            "Late-night exec checks, manager self-serve"),
              ("Governance carries through", "Sensitivity labels, RLS, audit log - same as PBI"),
            ], size=12)

# ----- SLIDE 8: RLS roles ---------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Row-Level Security", "Three roles, three persona views")


def role_card(x, y, w, h, title, sub, dax):
    add_rect(s, x, y, w, h, fill=LIGHTBG)
    add_text(s, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
             title, size=14, bold=True, color=INT_ORANGE)
    add_text(s, x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), Inches(0.6),
             sub, size=11, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.4), w - Inches(0.4), Inches(0.8),
             dax, size=10, color=GRAY)


role_card(Inches(0.5),  Inches(1.5), Inches(4.0), Inches(5.4),
          "HR Business Partner",
          "Reads all - org-wide HR analytics, executive view across functions and departments.",
          "modelPermission: read\n(no row filter)")
role_card(Inches(4.65), Inches(1.5), Inches(4.0), Inches(5.4),
          "People Manager",
          "Sees self + direct reports only - based on the dim_employee[manager_id] link to USERPRINCIPALNAME().",
          "tablePermission dim_employee =\n  [manager_id] = USERPRINCIPALNAME()\n  || [employee_id] = USERPRINCIPALNAME()")
role_card(Inches(8.80), Inches(1.5), Inches(4.0), Inches(5.4),
          "Compliance Officer",
          "Filters to employees in the Risk & Compliance function for COI and training oversight.",
          "tablePermission dim_employee =\n  RELATED(dim_department[function])\n  = \"Risk & Compliance\"")

# ----- SLIDE 9: Governance --------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Governance posture", "Built for OSFI E-21 expectations and Interac internal stewardship")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
            [
              ("Data residency",          "Canada-region Fabric capacity. All processing within Canadian boundary."),
              ("Sensitivity + Purview",   "Sensitivity labels propagate from source to dataset, report, and Copilot answers."),
              ("RLS at the model level",  "Persona filters enforced server-side, including Copilot and Data Agent surfaces."),
              ("Audit and lineage",       "Fabric lineage view + M365 audit log - dataset, measure, and Copilot usage all logged."),
              ("Responsible AI",          "Aligned with Microsoft Responsible AI Standard - groundedness, safety filters."),
              ("Change control",          "Semantic model definition is git-tracked (TMDL via API) - PR-style review of measure changes."),
            ], size=13)

# ----- SLIDE 10: Deployment + licensing -------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Deployment and licensing", "What's needed to take this live at Interac")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
            [
              ("Fabric capacity",         "F-SKU on Canada region. Demo runs on a trial - production sizing based on concurrency and data volume."),
              ("Power BI Pro / PPU",      "Pro for content authors and viewers; PPU if Direct Lake-on-SQL fallback is needed."),
              ("Copilot prerequisites",   "Per-user Copilot for Microsoft 365 licensing or workspace-scoped Copilot via Fabric capacity."),
              ("Identity",                "Entra ID groups for roles - HRBP, People Manager, Compliance Officer. SCIM provisioning ok."),
              ("Storage",                 "OneLake regional storage - same Canada geo as compute."),
              ("Networking",              "Private Link supported. Service tags available for IP egress controls."),
            ], size=13)

# ----- SLIDE 11: Roadmap ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Roadmap", "Where this stack is going")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
            [
              ("GA today",                 "Direct Lake on OneLake, Copilot in PBI, Smart Narrative, Find Insights, RLS"),
              ("Recent additions",         "Fabric Data Agent GA, custom example prompts on semantic models, multi-turn Copilot"),
              ("Near-term",                "Workspace-scoped Copilot governance, custom instructions per role, cross-report grounding"),
              ("Interac-specific opportunities", "Operational-risk Data Agent on OSFI E-21 metrics, Tier-1 risk scoring, vendor risk views"),
            ], size=13)

# ----- SLIDE 12: Next steps -------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Recommended next steps", "From demo to first pilot at Interac")
add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
            [
              ("Week 1",       "Identify pilot data domain (HR, Risk, Operations) and 1 executive sponsor"),
              ("Week 2",       "Stand up Fabric capacity in Canada region + provision pilot workspaces"),
              ("Week 3-4",     "Build first production semantic model with measure descriptions + RLS"),
              ("Week 5-6",     "Pilot users: 10-15 across HRBP / Manager / Compliance personas"),
              ("Week 7-8",     "Wire Copilot + Data Agent, capture prompt patterns, refine examplePrompts"),
              ("Week 9-12",    "Governance attestation (OSFI mapping) and broader rollout plan"),
              ("Microsoft Canada support", "Engineering pairing, weekly office hours, Frontier Transformation engagement on roadmap"),
            ], size=13)

# ----- SLIDE 13: Closing ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, Inches(0.6), Inches(3.0), Inches(1.0), Inches(0.10), fill=INT_ORANGE)
add_text(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.0),
         "Questions?", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.6),
         "Power BI + Fabric for Interac HR - what would you ask next?",
         size=18, color=RGBColor(0xCC, 0xD3, 0xDE))
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "anuragdh@microsoft.com  |  Microsoft Canada Data Specialist",
         size=11, color=RGBColor(0x9A, 0xA8, 0xBB))

prs.save(OUT)
print(f"Wrote {OUT}")
