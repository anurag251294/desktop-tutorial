"""Build the TCHC Fabric deep dive deck.

Slides carry the parts of the agenda that are discussion; the demo is the star and its
slides are section markers, not content. Nothing here duplicates what will be on screen
live -- a slide that competes with the demo just splits attention.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

sys.stdout.reconfigure(encoding="utf-8")

OUT = Path(r"C:/Users/ANURAG~1/AppData/Local/Temp/tchc/TCHC_Fabric_Deep_Dive.pptx")
SHOTS = Path(r"C:/Users/ANURAG~1/AppData/Local/Temp/tchc/shots")

# Civic palette, matching the report so the deck and the demo read as one thing.
INK = RGBColor(0x1F, 0x32, 0x43)
TEAL = RGBColor(0x2F, 0x6F, 0x6B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
AMBER = RGBColor(0xC0, 0x8A, 0x2E)
RED = RGBColor(0xB3, 0x42, 0x3A)
PAPER = RGBColor(0xF5, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    for index, run in enumerate(runs):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        if spacing:
            para.space_after = Pt(spacing)
        if run.get("bullet"):
            para.level = run.get("level", 0)
        piece = para.add_run()
        piece.text = run["t"]
        font = piece.font
        font.size = Pt(run.get("size", 16))
        font.bold = run.get("bold", False)
        font.color.rgb = run.get("color", INK)
        font.name = run.get("face", "Segoe UI")
    return box


def rule(s, x, y, w, color=TEAL, thickness=3):
    bar = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(thickness))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def panel(s, x, y, w, h, color=PAPER):
    box = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def header(s, eyebrow, title):
    text(s, 0.9, 0.55, 11.5, 0.35,
         [{"t": eyebrow.upper(), "size": 11, "bold": True, "color": TEAL,
           "face": "Consolas"}])
    text(s, 0.9, 0.95, 11.5, 0.8,
         [{"t": title, "size": 32, "color": INK, "face": "Georgia"}])
    rule(s, 0.9, 1.85, 1.2)


def bullets(s, x, y, w, items, size=15, gap=13):
    runs = []
    for item in items:
        if isinstance(item, tuple):
            body, level = item
        else:
            body, level = item, 0
        runs.append({"t": ("• " if level == 0 else "– ") + body,
                     "size": size if level == 0 else size - 1,
                     "color": INK if level == 0 else GRAY,
                     "level": level, "bullet": True})
    text(s, x, y, w, 4.5, runs, spacing=gap)


# ---------------------------------------------------------------- 1. title
s = slide(INK)
text(s, 1.0, 2.3, 11.3, 0.4,
     [{"t": "MICROSOFT FABRIC  ·  DEEP DIVE AND DEMO", "size": 13, "bold": True,
       "color": RGBColor(0x7F, 0xC4, 0xBE), "face": "Consolas"}])
text(s, 1.0, 2.85, 11.3, 1.5,
     [{"t": "Toronto Community Housing", "size": 46, "color": WHITE, "face": "Georgia"}])
text(s, 1.0, 4.05, 9.5, 1.0,
     [{"t": "Arrears and vacancy, end to end — from source extract to the decision a "
            "caseworker actually makes.", "size": 18,
       "color": RGBColor(0xB4, 0xC1, 0xC2), "face": "Georgia"}])
text(s, 1.0, 5.6, 11.3, 0.6,
     [{"t": "Thursday 27 August 2026   ·   2 hours   ·   Demonstration-led",
       "size": 13, "color": RGBColor(0x84, 0x96, 0x97), "face": "Consolas"}])

# ---------------------------------------------------------------- 2. objectives
s = slide()
header(s, "Item 1", "What we want you to leave with")
panel(s, 0.9, 2.25, 5.5, 3.6)
text(s, 1.25, 2.6, 4.9, 0.4,
     [{"t": "A WORKING PICTURE", "size": 11, "bold": True, "color": TEAL,
       "face": "Consolas"}])
bullets(s, 1.25, 3.05, 4.9, [
    "Fabric shown, not described",
    "Anchored to arrears and vacancy, not a generic retail demo",
    "Every number on screen traceable to a row",
])
panel(s, 6.9, 2.25, 5.5, 3.6)
text(s, 7.25, 2.6, 4.9, 0.4,
     [{"t": "A CONFIRMED LIST", "size": 11, "bold": True, "color": AMBER,
       "face": "Consolas"}])
bullets(s, 7.25, 3.05, 4.9, [
    "What has to be in place before kickoff",
    "An owner and a date against each item",
    "The kickoff date itself",
])
text(s, 0.9, 6.15, 11.5, 0.6,
     [{"t": "The technical and security questions should surface today, not in the "
            "first sprint.", "size": 15, "color": GRAY, "face": "Georgia"}])

# ---------------------------------------------------------------- 3. agenda
s = slide()
header(s, "Agenda", "Two hours")
rows = [
    ("1–2", "Objectives and where the engagement stands", "15 min", GRAY),
    ("3", "Fabric orientation — the platform in one picture", "10 min", GRAY),
    ("4", "Demo: landing TCHC data in OneLake", "20 min", TEAL),
    ("5", "Demo: medallion architecture in practice", "20 min", TEAL),
    ("6", "Demo: from curated data to a decision", "20 min", TEAL),
    ("7", "Security, privacy and data protection", "15 min", GRAY),
    ("8", "Governance, capacity and operations", "10 min", GRAY),
    ("9–10", "Open questions, and the pre-kickoff checklist", "10 min", AMBER),
]
y = 2.3
for number, label, mins, color in rows:
    text(s, 0.95, y, 0.9, 0.35,
         [{"t": number, "size": 13, "bold": True, "color": color, "face": "Consolas"}])
    text(s, 1.95, y, 8.4, 0.35, [{"t": label, "size": 15, "color": INK}])
    text(s, 10.6, y, 1.6, 0.35,
         [{"t": mins, "size": 12, "color": GRAY, "face": "Consolas"}])
    y += 0.52
text(s, 0.95, 6.65, 11.4, 0.4,
     [{"t": "Sixty of the hundred and twenty minutes are live demo. Questions are "
            "welcome throughout.", "size": 13, "color": GRAY, "face": "Georgia"}])

# ---------------------------------------------------------------- 4. platform
s = slide()
header(s, "Item 3", "The platform in one picture")
blocks = [
    ("OneLake", "One tenant-wide data lake.\nOne copy. Open Delta\nParquet.", TEAL),
    ("Lakehouse\nWarehouse", "Where the data lives\nand is shaped. Spark or\nT-SQL, same storage.", INK),
    ("Data Factory", "Pipelines and dataflows\nthat bring the data in\nand orchestrate it.", INK),
    ("Semantic model", "The shared definition of\na measure. Direct Lake,\nno second copy.", TEAL),
    ("Power BI", "Where a decision\nactually gets made.", INK),
]
x = 0.9
for title, body, color in blocks:
    panel(s, x, 2.35, 2.24, 2.5)
    rule(s, x, 2.35, 2.24, color, 4)
    text(s, x + 0.2, 2.6, 1.9, 0.7,
         [{"t": title, "size": 14, "bold": True, "color": color}])
    text(s, x + 0.2, 3.45, 1.9, 1.3, [{"t": body, "size": 11, "color": GRAY}])
    x += 2.36
panel(s, 0.9, 5.25, 11.5, 1.25, RGBColor(0xEC, 0xF0, 0xF0))
text(s, 1.25, 5.5, 10.8, 0.8,
     [{"t": "What this engagement uses: ", "size": 15, "bold": True, "color": INK},
      {"t": "all five. What it does not touch today — Real-Time Intelligence, "
            "Copilot and data agents, and migration of the existing Power BI estate.",
       "size": 15, "color": GRAY}])

# ---------------------------------------------------------------- 5. demo 1
s = slide(INK)
text(s, 1.0, 2.6, 11.3, 0.4,
     [{"t": "DEMO 1  ·  ITEM 4  ·  20 MINUTES", "size": 12, "bold": True,
       "color": RGBColor(0x7F, 0xC4, 0xBE), "face": "Consolas"}])
text(s, 1.0, 3.1, 11.3, 1.0,
     [{"t": "Landing TCHC data in OneLake", "size": 38, "color": WHITE,
       "face": "Georgia"}])
text(s, 1.0, 4.35, 10.5, 1.6,
     [{"t": "Shortcut, mirroring, pipeline, dataflow, gateway — which one fits which "
            "TCHC source system, and what each costs you.", "size": 17,
       "color": RGBColor(0xB4, 0xC1, 0xC2), "face": "Georgia"}])

# ---------------------------------------------------------------- 6. ingestion
s = slide()
header(s, "Item 4", "Five ways in, and what each is for")
cols = [("Option", 0.95, 2.6), ("Fits", 4.1, 5.6), ("Shown live", 10.3, 1.9)]
for label, x, w in cols:
    text(s, x, 2.3, w, 0.3,
         [{"t": label.upper(), "size": 10, "bold": True, "color": GRAY,
           "face": "Consolas"}])
rows = [
    ("Shortcut", "Data already in ADLS or another lakehouse. No copy, no latency.",
     "Yes", TEAL),
    ("Mirroring", "A supported source database, replicated continuously.",
     "No — needs a source", AMBER),
    ("Data Factory pipeline", "A system with an API or database, on a schedule.",
     "Yes", TEAL),
    ("Dataflow Gen2", "Low-code shaping, for a Power Query audience.",
     "Discussion", GRAY),
    ("On-prem / VNet gateway", "Whichever TCHC systems are not internet-reachable.",
     "No — needs a host", AMBER),
]
y = 2.75
for option, fits, live, color in rows:
    text(s, 0.95, y, 3.0, 0.4, [{"t": option, "size": 14, "bold": True, "color": INK}])
    text(s, 4.1, y, 6.0, 0.4, [{"t": fits, "size": 13, "color": GRAY}])
    text(s, 10.3, y, 2.2, 0.4,
         [{"t": live, "size": 12, "bold": True, "color": color, "face": "Consolas"}])
    y += 0.72
panel(s, 0.95, 6.35, 11.4, 0.85, RGBColor(0xFA, 0xF4, 0xE8))
text(s, 1.25, 6.55, 10.9, 0.5,
     [{"t": "Mirroring and the gateway need something on TCHC's side that does not "
            "exist yet. Architecture conversation today, not demo.",
       "size": 14, "color": INK, "face": "Georgia"}])

# ---------------------------------------------------------------- 7. demo 2
s = slide(INK)
text(s, 1.0, 2.6, 11.3, 0.4,
     [{"t": "DEMO 2  ·  ITEM 5  ·  20 MINUTES", "size": 12, "bold": True,
       "color": RGBColor(0x7F, 0xC4, 0xBE), "face": "Consolas"}])
text(s, 1.0, 3.1, 11.3, 1.0,
     [{"t": "The medallion, on live data", "size": 38, "color": WHITE,
       "face": "Georgia"}])
text(s, 1.0, 4.35, 10.5, 1.6,
     [{"t": "Bronze is what arrived, including its problems. Silver is what we are "
            "willing to stand behind. Gold is what a decision can rest on.",
       "size": 17, "color": RGBColor(0xB4, 0xC1, 0xC2), "face": "Georgia"}])

# ---------------------------------------------------------------- 8. silver
s = slide()
header(s, "Item 5", "What Silver actually has to fix")
pairs = [
    ("Two date formats", "dd/mm from the property system, mm/dd from finance"),
    ("Money as text", '"$1,234.56" arrives as a string'),
    ("Five spellings, two concepts", "RGI · rgi · Rent-Geared-to-Income · Market · MARKET"),
    ("A re-sent batch", "The extract ran twice; receipt_id deduplicates it"),
    ("Receipts for unknown accounts", "Quarantined, never dropped"),
    ("Work orders with no unit", "Quarantined, never dropped"),
]
y = 2.35
for left, right in pairs:
    text(s, 0.95, y, 4.3, 0.4, [{"t": left, "size": 14, "bold": True, "color": INK}])
    text(s, 5.4, y, 7.0, 0.4, [{"t": right, "size": 13, "color": GRAY,
                                "face": "Consolas"}])
    y += 0.6
panel(s, 0.95, 6.1, 11.4, 1.05, RGBColor(0xF7, 0xEE, 0xED))
text(s, 1.25, 6.28, 10.9, 0.75,
     [{"t": "Quarantine, not deletion.  ", "size": 15, "bold": True, "color": RED},
      {"t": "Silently dropping an orphan receipt makes arrears read higher than it is, "
            "and nobody downstream can tell it happened.", "size": 15, "color": INK,
       "face": "Georgia"}])

# ---------------------------------------------------------------- 9. demo 3
s = slide(INK)
text(s, 1.0, 2.6, 11.3, 0.4,
     [{"t": "DEMO 3  ·  ITEM 6  ·  20 MINUTES", "size": 12, "bold": True,
       "color": RGBColor(0x7F, 0xC4, 0xBE), "face": "Consolas"}])
text(s, 1.0, 3.1, 11.3, 1.0,
     [{"t": "From curated data to a decision", "size": 38, "color": WHITE,
       "face": "Georgia"}])
text(s, 1.0, 4.35, 10.5, 1.6,
     [{"t": "The closest view of what the MVP looks like. This is also where the "
            "modelling decisions start to matter.", "size": 17,
       "color": RGBColor(0xB4, 0xC1, 0xC2), "face": "Georgia"}])

# ---------------------------------------------------------------- 10. arrears calc
s = slide()
header(s, "Item 6", "Arrears is a calculation, not a column")
panel(s, 0.9, 2.3, 5.6, 2.5, RGBColor(0xF7, 0xEE, 0xED))
text(s, 1.25, 2.55, 5.0, 0.4,
     [{"t": "THE EASY VERSION", "size": 11, "bold": True, "color": RED,
       "face": "Consolas"}])
text(s, 1.25, 2.95, 5.0, 1.7,
     [{"t": "This month's charge minus this month's payment.\n\n"
            "A household that pays two weeks late every month looks permanently in "
            "arrears. One that missed a month last year and never made it up looks "
            "fine.", "size": 13, "color": INK}])
panel(s, 6.8, 2.3, 5.6, 2.5, RGBColor(0xF1, 0xF6, 0xF5))
text(s, 7.15, 2.55, 5.0, 0.4,
     [{"t": "WHAT WE DO", "size": 11, "bold": True, "color": TEAL, "face": "Consolas"}])
text(s, 7.15, 2.95, 5.0, 1.7,
     [{"t": "Receipts settle the oldest outstanding charge first, because that is how "
            "a rent account actually settles.\n\n"
            "The aging bucket follows the oldest charge still carrying a balance.",
       "size": 13, "color": INK}])
text(s, 0.9, 5.15, 11.5, 1.4,
     [{"t": "Both wrong answers reach a caseworker.", "size": 19, "bold": True,
       "color": INK, "face": "Georgia"},
      {"t": "One sends a collections letter to somebody who is current. The other "
            "misses a household that has been sliding for a year.", "size": 15,
       "color": GRAY, "face": "Georgia"}], spacing=10)

# ---------------------------------------------------------------- 11. direct lake
s = slide()
header(s, "Item 6", "Direct Lake, Import, DirectQuery")
modes = [
    ("Direct Lake", "Reads the Delta files directly.\nNo import. No refresh schedule.\nNo second copy.",
     "What we use", TEAL),
    ("Import", "A copy held in memory.\nFastest queries, but a refresh\nschedule and a latency.",
     "", GRAY),
    ("DirectQuery", "Every query goes to the source.\nAlways current, slowest, and it\nputs reporting load on the source.",
     "", GRAY),
]
x = 0.9
for title, body, tag, color in modes:
    panel(s, x, 2.35, 3.75, 2.9)
    rule(s, x, 2.35, 3.75, color, 4)
    text(s, x + 0.25, 2.6, 3.25, 0.45,
         [{"t": title, "size": 18, "bold": True, "color": color, "face": "Georgia"}])
    text(s, x + 0.25, 3.2, 3.25, 1.6, [{"t": body, "size": 13, "color": GRAY}])
    if tag:
        text(s, x + 0.25, 4.72, 3.25, 0.35,
             [{"t": tag.upper(), "size": 10, "bold": True, "color": color,
               "face": "Consolas"}])
    x += 3.93
panel(s, 0.9, 5.55, 11.5, 1.35, RGBColor(0xEC, 0xF0, 0xF0))
text(s, 1.25, 5.78, 10.8, 0.95,
     [{"t": "Direct Lake is why the medallion and the semantic layer are not two "
            "projects. ", "size": 15, "bold": True, "color": INK, "face": "Georgia"},
      {"t": "When the pipeline writes, the model sees it. The trade is one operational "
            "edge: replace a table rather than append to it and the model needs "
            "rebinding, not just refreshing.", "size": 15, "color": GRAY,
       "face": "Georgia"}])

# ---------------------------------------------------------------- 12. security
s = slide()
header(s, "Item 7", "Security, privacy and data protection")
left = [
    "Workspace roles — Admin, Member, Contributor, Viewer",
    "OneLake data access roles — folder level, so Gold without Bronze",
    "Row-level security — a caseworker sees their portfolio, a director sees all",
    "Object-level security — hide a column or a table from a role entirely",
]
right = [
    "Sensitivity labels via Purview, travelling with exports",
    "Entra ID groups as the single place access is granted",
    "Private endpoints — a design decision with your network team, early",
    "Tenant data in the dev environment — live, masked, or synthetic",
]
text(s, 0.95, 2.3, 5.3, 0.35,
     [{"t": "ACCESS", "size": 11, "bold": True, "color": TEAL, "face": "Consolas"}])
bullets(s, 0.95, 2.7, 5.3, left, size=14, gap=14)
text(s, 6.9, 2.3, 5.4, 0.35,
     [{"t": "PROTECTION AND NETWORK", "size": 11, "bold": True, "color": TEAL,
       "face": "Consolas"}])
bullets(s, 6.9, 2.7, 5.4, right, size=14, gap=14)
panel(s, 0.95, 5.85, 11.4, 1.15, RGBColor(0xFA, 0xF4, 0xE8))
text(s, 1.25, 6.05, 10.9, 0.85,
     [{"t": "A question back to you: ", "size": 15, "bold": True, "color": AMBER,
       "face": "Georgia"},
      {"t": "does the development environment use live, masked, or synthetic tenant "
            "data? Everything you will see today is synthetic. That is the safest "
            "default, and it should be a decision rather than an omission.",
       "size": 15, "color": INK, "face": "Georgia"}])

# ---------------------------------------------------------------- 13. governance
s = slide()
header(s, "Item 8", "Governance, capacity and operations")
items = [
    ("Domains and workspaces",
     "One workspace per environment per domain is the usual starting shape. "
     "Decide the naming convention before the first workspace, not the fifth."),
    ("Purview",
     "Lineage across OneLake, so \u201cwhere did this number come from\u201d has an "
     "answer that is not a person."),
    ("Capacity",
     "F-SKUs are a shared pool. Interactive queries and background jobs draw on the "
     "same units, and Fabric smooths bursts. The Capacity Metrics app shows what is "
     "actually consuming."),
    ("Git and deployment pipelines",
     "This entire demo is in a repository — notebooks, pipeline, model and report. "
     "Deployment pipelines promote between workspaces."),
]
y = 2.3
for title, body in items:
    text(s, 0.95, y, 3.4, 0.5, [{"t": title, "size": 15, "bold": True, "color": INK}])
    text(s, 4.5, y, 7.9, 0.9, [{"t": body, "size": 13, "color": GRAY}])
    y += 1.12
panel(s, 0.95, 6.55, 11.4, 0.72, RGBColor(0xEC, 0xF0, 0xF0))
text(s, 1.25, 6.7, 10.9, 0.45,
     [{"t": "We are not quoting an F-SKU today. It depends on volume, refresh pattern "
            "and concurrency, and none of those are known yet.", "size": 14,
       "color": INK, "face": "Georgia"}])

# ---------------------------------------------------------------- 13b. where AI fits
s = slide()
header(s, "Beyond this engagement", "Where the AI piece fits")
text(s, 0.95, 2.15, 11.4, 0.4,
     [{"t": "Out of scope for the MVP, and worth deciding deliberately rather than "
            "drifting into.", "size": 14, "color": GRAY, "face": "Georgia"}])

tiers = [
    ("ASK THE DATA", TEAL,
     "A Fabric data agent over the Gold layer.\n\n"
     "“Which buildings carry the most\nbalance over ninety days?” in plain\n"
     "language, against the same measures\nthe report uses.",
     "Closest to ready"),
    ("EXPLAIN THE DATA", TEAL,
     "A monthly arrears briefing, written\nfrom the pipeline's own numbers.\n\n"
     "The model never calculates. It renders\nwhat Gold produced, and every figure\n"
     "cites the row it came from.",
     "Built twice, elsewhere"),
    ("DECIDE WITH THE DATA", AMBER,
     "Ranking households for contact.\n\n"
     "This stops being a technology question.\nArrears touches somebody's housing,\n"
     "and that needs governance before it\nneeds a model.",
     "Not without a policy owner"),
]
x = 0.95
for label, colour, body, tag in tiers:
    panel(s, x, 2.65, 3.66, 3.0)
    rule(s, x, 2.65, 3.66, colour, 4)
    text(s, x + 0.25, 2.9, 3.2, 0.35,
         [{"t": label, "size": 11, "bold": True, "color": colour, "face": "Consolas"}])
    text(s, x + 0.25, 3.35, 3.2, 2.0, [{"t": body, "size": 12, "color": INK}])
    text(s, x + 0.25, 5.28, 3.2, 0.3,
         [{"t": tag, "size": 10, "bold": True, "color": GRAY, "face": "Consolas"}])
    x += 3.84

panel(s, 0.95, 5.9, 11.4, 1.35, RGBColor(0xF7, 0xEE, 0xED))
text(s, 1.25, 6.1, 10.9, 1.0,
     [{"t": "The guardrail that matters here.  ", "size": 15, "bold": True,
       "color": RED, "face": "Georgia"},
      {"t": "A household with a balance, a household we hold no data for, and a "
            "household outside the extract are three different things. Only the first "
            "is a finding. Any layer that lets them collapse into one list will "
            "eventually put the wrong name on it.", "size": 15, "color": INK,
       "face": "Georgia"}])

# ---------------------------------------------------------------- 13c. what it costs
s = slide()
header(s, "Beyond this engagement", "What the AI piece actually costs you")
notes = [
    ("Two tenant switches, not one",
     "A Fabric data agent needs Copilot and Azure AI enabled, and separately needs data "
     "permitted to leave the capacity's geography for processing. For an Ontario housing "
     "provider the second is a residency decision, not a checkbox."),
    ("Grounding is engineering, not prompting",
     "The pipeline decides what the model may see, and hands it a bounded evidence "
     "envelope. That work sits in the Gold layer and the contract between them — "
     "which is why it comes after the MVP, not instead of it."),
    ("Validate the contract, not just the values",
     "In an earlier build, an agent given no evidence returned a perfectly valid "
     "document citing a source that did not exist. Schema validation passed it. Only a "
     "second, independent check caught it."),
    ("A refusal has to be expressible",
     "If the only valid answer shape is an answer, a model will invent one. It needs a "
     "lawful way to return nothing."),
]
y = 2.3
for title, body in notes:
    text(s, 0.95, y, 3.6, 0.5, [{"t": title, "size": 15, "bold": True, "color": INK}])
    text(s, 4.7, y, 7.7, 0.95, [{"t": body, "size": 13, "color": GRAY}])
    y += 1.13
panel(s, 0.95, 6.65, 11.4, 0.72, RGBColor(0xEC, 0xF0, 0xF0))
text(s, 1.25, 6.8, 10.9, 0.45,
     [{"t": "None of this is a reason not to do it. It is the reason to do the "
            "medallion and the semantic model first.", "size": 14, "color": INK,
       "face": "Georgia"}])

# ---------------------------------------------------------------- 14. checklist
s = slide()
header(s, "Item 10", "Pre-kickoff checklist")
text(s, 0.95, 2.2, 11.4, 0.35,
     [{"t": "An owner and a date against each, before we leave the room.",
       "size": 14, "color": GRAY, "face": "Georgia"}])
checks = [
    "TCHC product owner named, with allocation confirmed for the engagement",
    "Fabric capacity decision — F-SKU size, reservation, Microsoft-hosted or TCHC tenant",
    "Azure subscription, Entra ID groups, administrative access for environment build",
    "Source system access for arrears and vacancy, plus a gateway host if required",
    "Azure DevOps or GitHub project available to the delivery team",
    "Development environment data — live, masked, or synthetic",
    "Security and privacy review slot booked inside the first two weeks",
]
y = 2.75
for item in checks:
    box = s.shapes.add_shape(1, Inches(0.98), Inches(y + 0.04), Inches(0.16),
                             Inches(0.16))
    box.fill.background()
    box.line.color.rgb = TEAL
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    text(s, 1.35, y - 0.04, 8.3, 0.4, [{"t": item, "size": 14, "color": INK}])
    text(s, 9.9, y - 0.04, 1.2, 0.4,
         [{"t": "owner", "size": 11, "color": RGBColor(0xC6, 0xCF, 0xCF),
           "face": "Consolas"}])
    text(s, 11.3, y - 0.04, 1.1, 0.4,
         [{"t": "date", "size": 11, "color": RGBColor(0xC6, 0xCF, 0xCF),
           "face": "Consolas"}])
    y += 0.56
rule(s, 0.95, 6.85, 11.4, TEAL, 2)
text(s, 0.95, 7.0, 11.4, 0.4,
     [{"t": "Kickoff: week of 14 September 2026", "size": 15, "bold": True,
       "color": INK, "face": "Georgia"}])

# ---------------------------------------------------------------- 15. close
s = slide(INK)
text(s, 1.0, 2.9, 11.3, 0.4,
     [{"t": "WHAT WE ARE ASKING FOR", "size": 12, "bold": True,
       "color": RGBColor(0x7F, 0xC4, 0xBE), "face": "Consolas"}])
text(s, 1.0, 3.4, 11.0, 1.4,
     [{"t": "The two answers that shape everything else", "size": 34, "color": WHITE,
       "face": "Georgia"}])
text(s, 1.0, 4.75, 10.6, 1.7,
     [{"t": "Which system is the system of record for arrears balances, and which for "
            "unit vacancy — and what integration surface each one offers.",
       "size": 17, "color": RGBColor(0xB4, 0xC1, 0xC2), "face": "Georgia"},
      {"t": "How current the data has to be for the decision the business actually "
            "makes. Daily, intraday, or near real time.", "size": 17,
       "color": RGBColor(0xB4, 0xC1, 0xC2), "face": "Georgia"}], spacing=14)

# ---------------------------------------------------------------- appendix shots
for index, (name, caption) in enumerate([
    ("final1.png", "Arrears overview — portfolio position, trend, aging, concentration"),
    ("final2.png", "Arrears deep dive — where the balance sits, and the caseload table"),
    ("final3.png", "Vacancy and turnaround — vacancy, revenue forgone, turnaround days"),
], start=1):
    path = SHOTS / name
    if not path.exists():
        continue
    s = slide()
    text(s, 0.9, 0.42, 11.5, 0.35,
         [{"t": f"APPENDIX  ·  REPORT PAGE {index}", "size": 11, "bold": True,
           "color": TEAL, "face": "Consolas"}])
    text(s, 0.9, 0.78, 11.5, 0.4, [{"t": caption, "size": 15, "color": INK,
                                    "face": "Georgia"}])
    s.shapes.add_picture(str(path), Inches(0.9), Inches(1.4), width=Inches(11.5))
    text(s, 0.9, 7.05, 11.5, 0.35,
         [{"t": "Synthetic data. 120 buildings, 10,568 units, 24 months.",
           "size": 11, "color": GRAY, "face": "Consolas"}])

prs.save(str(OUT))
print(f"wrote {OUT.name}  ({OUT.stat().st_size / 1024:.0f} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
