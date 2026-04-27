"""
Generate Manulife x Microsoft Fabric POC Presentation
Enterprise-grade PowerPoint for executive audiences
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x4E, 0x79)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = "Calibri"

BAR_HEIGHT = Inches(0.15)


def add_top_bar(slide):
    """Add thin navy bar at the top of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, BAR_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def set_tf(tf, text, size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    """Set text frame content."""
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=16, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4)):
    """Add a new paragraph to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_title_box(slide, text, left=Inches(0.7), top=Inches(0.4), width=Inches(11.5), height=Inches(0.8)):
    """Add a title text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    set_tf(tf, text, size=30, bold=True, color=NAVY)
    return txBox


def add_body_box(slide, left=Inches(0.7), top=Inches(1.4), width=Inches(11.5), height=Inches(5.0)):
    """Add a body text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_colored_rect(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE, bold=False):
    """Add a colored rectangle with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.name = FONT_NAME
    p.font.bold = bold
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    try:
        shape.text_frame.auto_size = None
    except Exception:
        pass
    return shape


def add_arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.3)):
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    return shape


def set_cell(cell, text, size=12, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    """Set table cell text."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def style_header_row(table, col_count):
    """Style the first row as a header."""
    for i in range(col_count):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.bold = True


def add_table(slide, rows_data, left, top, width, col_widths, row_height=Inches(0.45)):
    """Add a table to a slide. rows_data is list of lists of strings. First row is header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.45 * n_rows))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            set_cell(table.cell(r, c), val, size=13, bold=(r == 0),
                     color=WHITE if r == 0 else DARK_TEXT)
    style_header_row(table, n_cols)
    # Alternate row shading
    for r in range(2, n_rows, 2):
        for c in range(n_cols):
            table.cell(r, c).fill.solid()
            table.cell(r, c).fill.fore_color.rgb = LIGHT_GRAY
    return table


# ── Build Presentation ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
# Full navy background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Teal accent bar
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.0), SLIDE_W, Inches(0.06)).fill.solid()
slide.shapes.placeholders  # force refresh
accent = slide.shapes[-1]
accent.fill.solid()
accent.fill.fore_color.rgb = TEAL
accent.line.fill.background()

# Title text
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11), Inches(1.2))
tf = txBox.text_frame
set_tf(tf, "Manulife x Microsoft Fabric", size=44, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11), Inches(0.7))
tf2 = txBox2.text_frame
set_tf(tf2, "Intelligent Data & Analytics POC", size=26, bold=False, color=RGBColor(0xBB, 0xDD, 0xEE), alignment=PP_ALIGN.LEFT)

txBox3 = slide.shapes.add_textbox(Inches(1.0), Inches(3.4), Inches(11), Inches(0.7))
tf3 = txBox3.text_frame
set_tf(tf3, "Copilot-Style Experience with Fabric Data Agent", size=20, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT)

txBox4 = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(5), Inches(0.5))
tf4 = txBox4.text_frame
set_tf(tf4, "April 2026", size=16, bold=False, color=RGBColor(0x99, 0xBB, 0xCC), alignment=PP_ALIGN.LEFT)

txBox5 = slide.shapes.add_textbox(Inches(9.5), Inches(6.5), Inches(3), Inches(0.5))
tf5 = txBox5.text_frame
set_tf(tf5, "CONFIDENTIAL", size=14, bold=True, color=RGBColor(0xFF, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Agenda")

agenda_items = [
    "Business Context & Opportunity",
    "Solution Architecture",
    "Data Foundation (OneLake)",
    "Analytics Layer (Semantic Model)",
    "AI-Powered Access (Data Agent)",
    "Live Demo",
    "What's Feasible Now vs. Roadmap",
    "Recommended Next Steps",
]
tf = add_body_box(slide, top=Inches(1.5), height=Inches(5.5))
for i, item in enumerate(agenda_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {i+1}.   {item}"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_TEXT
    p.font.name = FONT_NAME
    p.space_before = Pt(10)
    p.space_after = Pt(6)
    # Color the number
    # (python-pptx doesn't easily let us color part of text without runs, keeping simple)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Business Challenge
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "The Opportunity")

bullets = [
    "Service reps and analysts need fast answers across policy, claims, and investment data",
    "Information is scattered across structured databases and unstructured documents",
    "Traditional BI dashboards require training and don't answer ad-hoc questions",
    "Need: A natural-language interface that combines data insights with document context",
    "Goal: Reduce time-to-insight from hours to seconds",
]
tf = add_body_box(slide, top=Inches(1.6), height=Inches(5.0))
for i, b in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = b
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.font.name = FONT_NAME
    p.space_before = Pt(12)
    p.space_after = Pt(8)
    p.level = 0
    # Add bullet character
    p.text = "\u2022  " + b


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4: POC Scope
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "POC Scope")

boxes_data = [
    ("Unified Data\nFoundation", "OneLake as single\nsource of truth"),
    ("Trusted\nBusiness Logic", "Semantic model with\nDAX measures & KPIs"),
    ("Natural Language\nAccess", "Fabric Data Agent for\nconversational queries"),
    ("Document\nEnrichment", "Policy docs, guidelines\nintegrated via RAG"),
]

for i, (title, desc) in enumerate(boxes_data):
    left = Inches(0.7 + i * 3.1)
    top = Inches(2.0)
    # Title box
    add_colored_rect(slide, left, top, Inches(2.8), Inches(1.2), NAVY, title, font_size=16, font_color=WHITE, bold=True)
    # Description box
    add_colored_rect(slide, left, top + Inches(1.2), Inches(2.8), Inches(1.5), LIGHT_GRAY, desc, font_size=14, font_color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Architecture Philosophy
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Architecture Philosophy")

# Key message
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "OneLake is the foundation.  The Semantic Model is the brain.  The Data Agent is the voice."
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = TEAL
p.font.name = FONT_NAME
p.alignment = PP_ALIGN.CENTER

# Three-layer visual
layers = [
    ("Users  \u2192  Copilot Experience", TEAL, Inches(2.8)),
    ("Data Agent  \u2192  Semantic Model", NAVY, Inches(3.8)),
    ("OneLake  (Bronze \u2192 Silver \u2192 Gold)", RGBColor(0x3A, 0x6B, 0x35), Inches(4.8)),
]
for text, color, top in layers:
    add_colored_rect(slide, Inches(2.0), top, Inches(7.5), Inches(0.7), color, text, font_size=16, font_color=WHITE, bold=True)

# Side enrichment
add_colored_rect(slide, Inches(10.0), Inches(3.2), Inches(2.8), Inches(2.5), RGBColor(0x8B, 0x5C, 0x2A),
                 "Unstructured Docs\n(Azure AI Search)\nas enrichment", font_size=13, font_color=WHITE, bold=False)

# Arrows between layers
for top in [Inches(3.55), Inches(4.55)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(5.7), top, Inches(0.3), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_GRAY
    shape.line.fill.background()


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6: End-to-End Architecture
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Reference Architecture")

# Build architecture diagram with shapes
cx = Inches(5.0)  # center x for main column

# Row 1 - Users
add_colored_rect(slide, Inches(4.5), Inches(1.3), Inches(4.0), Inches(0.55), TEAL, "Business Users", font_size=14, font_color=WHITE, bold=True)

# Arrow
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(1.85), Inches(0.25), Inches(0.2)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = MED_GRAY
slide.shapes[-1].line.fill.background()

# Row 2 - Copilot
add_colored_rect(slide, Inches(4.5), Inches(2.1), Inches(4.0), Inches(0.55), NAVY, "Standalone Copilot Experience", font_size=13, font_color=WHITE, bold=True)

slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(2.65), Inches(0.25), Inches(0.2)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = MED_GRAY
slide.shapes[-1].line.fill.background()

# Row 3 - Orchestration + Data Agent
add_colored_rect(slide, Inches(3.0), Inches(2.9), Inches(3.0), Inches(0.55), RGBColor(0x4A, 0x7C, 0x59), "Azure OpenAI Orchestration", font_size=12, font_color=WHITE, bold=True)
add_colored_rect(slide, Inches(6.2), Inches(2.9), Inches(3.0), Inches(0.55), TEAL, "Fabric Data Agent", font_size=13, font_color=WHITE, bold=True)

# Row 4 - Semantic Model
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.5), Inches(3.45), Inches(0.25), Inches(0.2)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = MED_GRAY
slide.shapes[-1].line.fill.background()

add_colored_rect(slide, Inches(4.5), Inches(3.7), Inches(4.0), Inches(0.55), NAVY, "Semantic Model (Power BI)", font_size=13, font_color=WHITE, bold=True)

# Row 5 - OneLake
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(4.25), Inches(0.25), Inches(0.2)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = MED_GRAY
slide.shapes[-1].line.fill.background()

add_colored_rect(slide, Inches(3.5), Inches(4.5), Inches(2.0), Inches(0.55), RGBColor(0xB8, 0x86, 0x0B), "Bronze (Raw)", font_size=12, font_color=WHITE, bold=True)
add_colored_rect(slide, Inches(5.6), Inches(4.5), Inches(2.0), Inches(0.55), RGBColor(0xA0, 0xA0, 0xA0), "Silver (Clean)", font_size=12, font_color=WHITE, bold=True)
add_colored_rect(slide, Inches(7.7), Inches(4.5), Inches(2.0), Inches(0.55), RGBColor(0xDA, 0xA5, 0x20), "Gold (Star)", font_size=12, font_color=WHITE, bold=True)

# OneLake label
add_colored_rect(slide, Inches(3.5), Inches(5.1), Inches(6.2), Inches(0.35), LIGHT_GRAY, "OneLake Lakehouse", font_size=12, font_color=DARK_TEXT, bold=True)

# Data sources at bottom
add_colored_rect(slide, Inches(2.0), Inches(5.8), Inches(3.5), Inches(0.55), RGBColor(0x5B, 0x7E, 0x91), "Data Sources: CSV / DB (Structured)", font_size=11, font_color=WHITE)
add_colored_rect(slide, Inches(5.8), Inches(5.8), Inches(3.5), Inches(0.55), RGBColor(0x5B, 0x7E, 0x91), "Documents: PDFs / Docs (Unstructured)", font_size=11, font_color=WHITE)

# Azure AI Search (side)
add_colored_rect(slide, Inches(10.5), Inches(2.9), Inches(2.3), Inches(1.0), RGBColor(0x8B, 0x5C, 0x2A), "Azure AI Search\n(Doc Retrieval)", font_size=12, font_color=WHITE, bold=True)

# Governance sidebar
add_colored_rect(slide, Inches(0.3), Inches(1.3), Inches(2.0), Inches(5.3), RGBColor(0x44, 0x44, 0x44), "Governance\n&\nSecurity\n\n\u2022 RLS\n\u2022 RBAC\n\u2022 Audit\n\u2022 Lineage", font_size=11, font_color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Data Flow — Structured
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Data Flow \u2014 Structured Path")

flow_items = [
    ("Raw\nSources", RGBColor(0x5B, 0x7E, 0x91)),
    ("Fabric\nPipelines", TEAL),
    ("Bronze\n(Raw)", RGBColor(0xB8, 0x86, 0x0B)),
    ("Notebooks\n(Transform)", RGBColor(0x4A, 0x7C, 0x59)),
    ("Silver\n(Clean)", RGBColor(0xA0, 0xA0, 0xA0)),
    ("Gold\n(Star Schema)", RGBColor(0xDA, 0xA5, 0x20)),
    ("Semantic\nModel", NAVY),
    ("Data\nAgent", TEAL),
    ("User", NAVY),
]

box_w = Inches(1.2)
box_h = Inches(0.9)
start_x = Inches(0.5)
y = Inches(3.0)
spacing = Inches(1.4)

for i, (label, color) in enumerate(flow_items):
    x = start_x + i * spacing
    add_colored_rect(slide, x, y, box_w, box_h, color, label, font_size=11, font_color=WHITE, bold=True)
    if i < len(flow_items) - 1:
        add_arrow_right(slide, x + box_w + Inches(0.02), y + Inches(0.3), width=Inches(0.15), height=Inches(0.25))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Data Flow — Unstructured
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Data Flow \u2014 Unstructured Enrichment")

flow2 = [
    ("Documents\n(PDFs, Guidelines,\nReports)", RGBColor(0x5B, 0x7E, 0x91)),
    ("Document\nProcessing\n(Chunking)", TEAL),
    ("Azure AI Search\n(Vector Index)", RGBColor(0x8B, 0x5C, 0x2A)),
    ("Orchestration\nLayer", RGBColor(0x4A, 0x7C, 0x59)),
    ("Combined\nResponse", NAVY),
]

box_w2 = Inches(2.0)
box_h2 = Inches(1.1)
start_x2 = Inches(0.7)
y2 = Inches(2.5)
spacing2 = Inches(2.5)

for i, (label, color) in enumerate(flow2):
    x = start_x2 + i * spacing2
    add_colored_rect(slide, x, y2, box_w2, box_h2, color, label, font_size=12, font_color=WHITE, bold=True)
    if i < len(flow2) - 1:
        add_arrow_right(slide, x + box_w2 + Inches(0.05), y2 + Inches(0.4), width=Inches(0.35), height=Inches(0.25))

# Note about merging
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(11.5), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
set_tf(tf, "Parallel path merges at orchestration layer:", size=16, bold=True, color=NAVY)
add_paragraph(tf, "\u2022  Structured data answers come from Semantic Model via Data Agent", size=15, color=DARK_TEXT)
add_paragraph(tf, "\u2022  Unstructured context comes from Azure AI Search via RAG pattern", size=15, color=DARK_TEXT)
add_paragraph(tf, "\u2022  Orchestration combines both into a unified, grounded response", size=15, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Data Foundation — OneLake
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "OneLake \u2014 The Data Foundation")

table_data = [
    ["Layer", "Purpose", "Tables", "Format"],
    ["Bronze", "Raw ingestion", "7 tables", "Delta"],
    ["Silver", "Cleansed & enriched", "7 tables", "Delta"],
    ["Gold", "Star schema (facts + dims)", "10 tables", "Delta"],
    ["Documents", "Chunked content", "1 table", "Delta"],
]
add_table(slide, table_data, Inches(1.5), Inches(1.8), Inches(10.0),
          [Inches(2.0), Inches(4.0), Inches(2.0), Inches(2.0)])

# Summary
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.0), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
set_tf(tf, "Total: 26 delta tables  |  339 KB raw data  |  Open Delta Parquet format", size=17, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Gold Layer — Star Schema
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Curated Analytics Layer")

# Center: Fact tables
fact_tables = ["fact_claims", "fact_transactions", "fact_investments", "fact_policy_premiums"]
for i, t in enumerate(fact_tables):
    add_colored_rect(slide, Inches(4.5), Inches(2.0 + i * 0.75), Inches(3.0), Inches(0.6), NAVY, t, font_size=13, font_color=WHITE, bold=True)

# Dimension tables around
dim_tables_left = ["dim_customer", "dim_product", "dim_advisor"]
dim_tables_right = ["dim_policy", "dim_date", "dim_fund"]

for i, t in enumerate(dim_tables_left):
    add_colored_rect(slide, Inches(1.0), Inches(2.0 + i * 1.0), Inches(2.5), Inches(0.6), TEAL, t, font_size=13, font_color=WHITE, bold=True)

for i, t in enumerate(dim_tables_right):
    add_colored_rect(slide, Inches(8.5), Inches(2.0 + i * 1.0), Inches(2.5), Inches(0.6), TEAL, t, font_size=13, font_color=WHITE, bold=True)

# Lines connecting (use thin rectangles as connectors)
for i in range(3):
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(2.25 + i * 1.0), Inches(1.0), Inches(0.02)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = MED_GRAY
    slide.shapes[-1].line.fill.background()
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(2.25 + i * 1.0), Inches(1.0), Inches(0.02)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = MED_GRAY
    slide.shapes[-1].line.fill.background()

# Legend
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(10.0), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
set_tf(tf, "Star schema design: 4 fact tables + 6 dimension tables", size=16, bold=False, color=DARK_TEXT)
add_paragraph(tf, "Key relationships: Customer \u2192 Policy \u2192 Claims/Premiums | Customer \u2192 Investments | Advisor \u2192 Customer | Product \u2192 Policy", size=14, color=RGBColor(0x66, 0x66, 0x66))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Sample Data Overview
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "POC Dataset")

table_data = [
    ["Dataset", "Records", "Description"],
    ["Customers", "200", "Canadian insurance customers"],
    ["Policies", "391", "Life, Health, Auto, Home, Travel, Disability"],
    ["Claims", "300", "~70% approved, ~20% denied"],
    ["Products", "25", "Insurance, Investment, Annuity"],
    ["Investments", "300", "Equity, Bond, Balanced, Money Market, Real Estate"],
    ["Advisors", "30", "6 Canadian regions"],
    ["Transactions", "800", "Payments, payouts, purchases"],
    ["Documents", "8", "Policy terms, guidelines, FAQs, reports"],
]
add_table(slide, table_data, Inches(1.5), Inches(1.6), Inches(10.0),
          [Inches(2.5), Inches(1.5), Inches(6.0)])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Semantic Model
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Semantic Model \u2014 Trusted KPIs")

# Left column - DAX measures
add_colored_rect(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.5), NAVY, "14 DAX Measures", font_size=15, font_color=WHITE, bold=True)

measures = [
    "Total Premium Revenue",
    "Total Claims Amount / Approved Amount",
    "Claim Count / Approval Rate",
    "Average Processing Days",
    "Total AUM / Investment Inflows",
    "Average Return YTD",
    "Policy Count / Total Coverage",
    "Transaction Amount / Count",
]
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, m in enumerate(measures):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"\u2022  {m}"
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_TEXT
    p.font.name = FONT_NAME
    p.space_before = Pt(6)

# Right column - Why Semantic Model Matters
add_colored_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5), TEAL, "Why Semantic Model Matters", font_size=15, font_color=WHITE, bold=True)

reasons = [
    "Single source of truth for business logic",
    "Consistent calculations across all consumers",
    "Data Agent queries measures, not raw tables",
    "Governance: who can see what",
]
txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, r in enumerate(reasons):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = f"\u2022  {r}"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.font.name = FONT_NAME
    p.space_before = Pt(10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Fabric Data Agent
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Natural Language Access \u2014 Fabric Data Agent")

# Key points
key_points = [
    "Queries the semantic model using natural language",
    "Returns structured answers with data visualizations",
    "No SQL or DAX knowledge required",
    "Governed: respects row-level security",
    "Currently in Preview \u2014 GA expected 2025-2026",
]
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(6.5), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True
for i, pt in enumerate(key_points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"\u2022  {pt}"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.font.name = FONT_NAME
    p.space_before = Pt(8)

# Sample questions box
add_colored_rect(slide, Inches(7.5), Inches(1.5), Inches(5.3), Inches(0.5), NAVY, "Sample Questions", font_size=14, font_color=WHITE, bold=True)

questions = [
    '"How many active policies do we have?"',
    '"Show claims ratio by policy type"',
    '"Top 5 customers by investment AUM"',
    '"What\'s the premium trend over last 6 months?"',
]
txBox2 = slide.shapes.add_textbox(Inches(7.5), Inches(2.2), Inches(5.3), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, q in enumerate(questions):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = q
    p.font.size = Pt(14)
    p.font.color.rgb = TEAL
    p.font.name = FONT_NAME
    p.font.italic = True
    p.space_before = Pt(10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Unstructured Enrichment
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Document Intelligence \u2014 Enrichment Layer")

points = [
    "8 enterprise documents processed and chunked",
    "Document types: policy terms, claims guidelines, product guides, FAQs, investment commentary, compliance handbooks",
    "Chunking: sliding window (500 tokens, 50 overlap)",
    "Pattern: Azure AI Search + Azure OpenAI for RAG",
]
tf = add_body_box(slide, top=Inches(1.5), width=Inches(11.5), height=Inches(2.5))
for i, pt in enumerate(points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"\u2022  {pt}"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.font.name = FONT_NAME
    p.space_before = Pt(8)

# Example box
add_colored_rect(slide, Inches(0.7), Inches(4.3), Inches(11.5), Inches(0.45), TEAL, "Example", font_size=14, font_color=WHITE, bold=True)

txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(11.5), Inches(1.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
set_tf(tf2, 'Q: "What does the claims guideline say about fraud indicators?"', size=15, bold=False, color=NAVY)
add_paragraph(tf2, "\u2192  Retrieves relevant document chunks from Azure AI Search", size=14, color=DARK_TEXT)
add_paragraph(tf2, "\u2192  Generates answer with document citations via Azure OpenAI", size=14, color=DARK_TEXT)
add_paragraph(tf2, "\u2192  Response is grounded in actual enterprise documents", size=14, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15: Live Demo (Section Divider)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
# Full navy background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Teal accent
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(4.0), Inches(7.0), Inches(0.05))
slide.shapes[-1].fill.solid()
slide.shapes[-1].fill.fore_color.rgb = TEAL
slide.shapes[-1].line.fill.background()

txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.0), Inches(1.2))
tf = txBox.text_frame
set_tf(tf, "Live Demo", size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(4.3), Inches(9.0), Inches(0.8))
tf2 = txBox2.text_frame
set_tf(tf2, "Fabric Data Agent in Action", size=24, bold=False, color=RGBColor(0xBB, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16: Demo Questions
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Demo Walkthrough")

demo_qs = [
    ("Simple", '"How many active customers do we have?"'),
    ("Analytical", '"Show claims by policy type and region"'),
    ("Comparative", '"Which advisors have the highest AUM?"'),
    ("Trend", '"What\'s the monthly premium trend?"'),
    ("Hybrid", '"Summarize claims data and relevant policy guidelines for Health insurance"'),
]

for i, (label, question) in enumerate(demo_qs):
    top = Inches(1.6 + i * 1.05)
    # Number circle
    add_colored_rect(slide, Inches(0.7), top, Inches(0.5), Inches(0.5), NAVY, str(i + 1), font_size=16, font_color=WHITE, bold=True)
    # Label
    add_colored_rect(slide, Inches(1.4), top, Inches(2.0), Inches(0.5), TEAL, label, font_size=14, font_color=WHITE, bold=True)
    # Question
    txBox = slide.shapes.add_textbox(Inches(3.7), top, Inches(8.5), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_tf(tf, question, size=15, bold=False, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17: What's Feasible Now
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Current Capabilities vs. Roadmap")

table_data = [
    ["Capability", "Status", "Notes"],
    ["OneLake + Medallion Architecture", "GA", "Fully available"],
    ["PySpark Notebooks", "GA", "Fully available"],
    ["Semantic Model (DirectLake)", "GA", "Fully available"],
    ["Fabric Data Agent", "Preview", "Limited GA features"],
    ["Native RAG in Data Agent", "Roadmap", "Requires Azure AI Search sidecar"],
    ["Multi-source grounding", "Roadmap", "Agent + docs in single query"],
    ["Copilot in Fabric", "GA (tenant opt-in)", "Admin must enable"],
]
add_table(slide, table_data, Inches(0.7), Inches(1.6), Inches(12.0),
          [Inches(4.5), Inches(2.5), Inches(5.0)])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18: Key Outcomes
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "POC Outcomes")

outcomes = [
    "End-to-end data pipeline from raw to Gold layer",
    "26 delta tables across medallion architecture",
    "Semantic model with 14 DAX measures and 6 relationships",
    "Document chunking pipeline for 8 enterprise documents",
    "Natural-language query capability via Data Agent",
    "Reference architecture suitable for production scale",
]

for i, outcome in enumerate(outcomes):
    top = Inches(1.7 + i * 0.85)
    # Checkmark box
    add_colored_rect(slide, Inches(0.7), top, Inches(0.55), Inches(0.55), RGBColor(0x2D, 0x8B, 0x57), "\u2713", font_size=20, font_color=WHITE, bold=True)
    # Text
    txBox = slide.shapes.add_textbox(Inches(1.5), top, Inches(10.5), Inches(0.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_tf(tf, outcome, size=18, bold=False, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19: Recommended Next Steps
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Recommended Next Steps")

table_data = [
    ["#", "Action", "Timeline"],
    ["1", "Enable tenant settings for Data Agent", "Week 1"],
    ["2", "Validate POC results with business stakeholders", "Week 2"],
    ["3", "Define production pilot scope (single BU)", "Week 3-4"],
    ["4", "Provision production Fabric capacity", "Week 4-5"],
    ["5", "Engage Microsoft FastTrack / co-sell", "Week 2-3"],
    ["6", "Develop production SOW", "Week 5-6"],
]
add_table(slide, table_data, Inches(1.0), Inches(1.6), Inches(11.0),
          [Inches(0.8), Inches(7.2), Inches(3.0)])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20: Investment Overview
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Investment Overview")

table_data = [
    ["Phase", "Estimate", "Timeline"],
    ["POC (current)", "$150K - $200K", "3-4 months"],
    ["Year 1 Production", "$1M - $2M", "12 months"],
    ["Steady State (Year 2+)", "$600K - $1.5M / yr", "Ongoing"],
    ["3-Year TCV (base case)", "$3M - $6M", ""],
]
add_table(slide, table_data, Inches(1.5), Inches(1.8), Inches(10.0),
          [Inches(4.0), Inches(3.0), Inches(3.0)])

# Note
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.0), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
set_tf(tf, "Note: Estimates are directional and subject to scoping. Includes Fabric capacity, development, and ongoing support.", size=13, bold=False, color=RGBColor(0x88, 0x88, 0x88))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 21: Why Microsoft Fabric
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_top_bar(slide)
add_title_box(slide, "Why Fabric for Manulife")

differentiators = [
    ("Unified Platform", "No stitching BI + data engineering + AI \u2014 one product, one experience"),
    ("OneLake", "One copy of data, open Delta format, governed from day one"),
    ("DirectLake", "Semantic model reads delta tables directly \u2014 no import, no duplication"),
    ("Copilot + Data Agent", "Natural language access as a first-class feature, not an afterthought"),
]

for i, (title, desc) in enumerate(differentiators):
    top = Inches(1.7 + i * 1.3)
    # Number
    add_colored_rect(slide, Inches(0.7), top, Inches(0.6), Inches(0.6), NAVY, str(i + 1), font_size=20, font_color=WHITE, bold=True)
    # Title
    txBox = slide.shapes.add_textbox(Inches(1.6), top - Inches(0.05), Inches(10.5), Inches(0.45))
    tf = txBox.text_frame
    set_tf(tf, title, size=20, bold=True, color=NAVY)
    # Description
    txBox2 = slide.shapes.add_textbox(Inches(1.6), top + Inches(0.4), Inches(10.5), Inches(0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    set_tf(tf2, desc, size=16, bold=False, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 22: Thank You
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
# Full navy background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Teal accent line
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.8), Inches(6.0), Inches(0.05))
slide.shapes[-1].fill.solid()
slide.shapes[-1].fill.fore_color.rgb = TEAL
slide.shapes[-1].line.fill.background()

txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.0), Inches(1.2))
tf = txBox.text_frame
set_tf(tf, "Thank You", size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(4.0), Inches(9.0), Inches(0.8))
tf2 = txBox2.text_frame
set_tf(tf2, "Questions & Discussion", size=24, bold=False, color=RGBColor(0xBB, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

txBox3 = slide.shapes.add_textbox(Inches(2.0), Inches(5.5), Inches(9.0), Inches(0.5))
tf3 = txBox3.text_frame
set_tf(tf3, "CONFIDENTIAL \u2014 Manulife POC Engagement", size=14, bold=True, color=RGBColor(0xFF, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────
output_dir = os.path.dirname(os.path.abspath(__file__)).replace("scripts", "docs")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "Manulife-Fabric-POC-Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
