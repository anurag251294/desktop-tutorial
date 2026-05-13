"""Build the Manulife Japan Fabric POC slide deck.

Audience: MLJ data leadership (L100-ish technical level). Centerpiece is the
side-by-side comparison: current ADF + Databricks + Unity Catalog stack vs.
Fabric, mapped to the layers in MLJ-Data-Architecture.pdf.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand colors
FABRIC_GREEN = RGBColor(0x11, 0x78, 0x65)
FABRIC_GREEN_LIGHT = RGBColor(0xC9, 0xE6, 0xDF)
DATABRICKS_RED = RGBColor(0xFF, 0x36, 0x21)
DATABRICKS_RED_LIGHT = RGBColor(0xFC, 0xDA, 0xD3)
GREY_DARK = RGBColor(0x33, 0x33, 0x33)
GREY_MID = RGBColor(0x59, 0x59, 0x59)
GREY_LIGHT = RGBColor(0xE8, 0xE8, 0xE8)
GREY_BG = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF2, 0xC8, 0x11)
INK = RGBColor(0x1F, 0x1F, 0x1F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]

# ------------------ helpers ------------------

def add_rect(slide, x, y, w, h, fill, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_width is not None:
            shp.line.width = line_width
    shp.shadow.inherit = False
    return shp

def add_round(slide, x, y, w, h, fill, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.12
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_width is not None:
            shp.line.width = line_width
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb

def add_bullets(slide, x, y, w, h, bullets, *, size=14, color=INK, bullet_color=None, bold_first_word=False, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        bullet_run = p.add_run()
        bullet_run.text = "•  "
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = bullet_color or color
        bullet_run.font.bold = True
        bullet_run.font.name = font
        if bold_first_word and " " in line:
            head, tail = line.split(" ", 1)
            r1 = p.add_run(); r1.text = head + " "
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name = font
            r2 = p.add_run(); r2.text = tail
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = font
        else:
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = font
    return tb

def slide_header(slide, title, subtitle=None, footer_idx=None):
    add_rect(slide, 0, 0, SW, Inches(0.06), FABRIC_GREEN)
    add_text(slide, Inches(0.5), Inches(0.18), SW - Inches(1.0), Inches(0.55),
             title, size=26, bold=True, color=INK)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.72), SW - Inches(1.0), Inches(0.35),
                 subtitle, size=13, color=GREY_MID)
    # Footer bar
    add_rect(slide, 0, SH - Inches(0.30), SW, Inches(0.30), GREY_BG)
    add_text(slide, Inches(0.5), SH - Inches(0.27), SW - Inches(2.0), Inches(0.24),
             "Manulife Japan x Microsoft  |  Fabric POC", size=9, color=GREY_MID)
    if footer_idx is not None:
        add_text(slide, SW - Inches(1.0), SH - Inches(0.27), Inches(0.6), Inches(0.24),
                 str(footer_idx), size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)

# ============== SLIDE 1: COVER ==============
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, FABRIC_GREEN)
# accent stripe
add_rect(s, 0, Inches(4.6), SW, Inches(0.04), ACCENT)
add_text(s, Inches(0.7), Inches(1.8), SW - Inches(1.4), Inches(0.6),
         "Manulife Japan", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(2.45), SW - Inches(1.4), Inches(1.0),
         "From Unity Catalog to Microsoft Fabric", size=32, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.4), SW - Inches(1.4), Inches(0.6),
         "A unified data + AI platform for the 2026 governance refresh", size=18, color=WHITE)
add_text(s, Inches(0.7), Inches(4.8), SW - Inches(1.4), Inches(0.5),
         "Technical merits, side-by-side comparison, and POC walkthrough", size=14, color=WHITE)
add_text(s, Inches(0.7), Inches(6.6), SW - Inches(1.4), Inches(0.4),
         "Microsoft  |  Prepared for Manulife Japan Data & AI Workshop", size=11, color=WHITE)

# ============== SLIDE 2: AGENDA ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Agenda", footer_idx=2)
items = [
    ("01", "Manulife Japan today", "Current architecture, what's working, where the friction is"),
    ("02", "Why Fabric for MLJ", "The technical case in one slide"),
    ("03", "Side-by-side comparison", "Current stack vs Fabric, layer by layer"),
    ("04", "OneLake + Unity Catalog story", "What changes for the 2026 governance refresh"),
    ("05", "AI-driven self-service", "Data Agent, Copilot, and self-service reporting"),
    ("06", "Live demo", "Walkthrough of the Fabric POC built for MLJ"),
    ("07", "Migration phasing", "Pragmatic path from current to target state"),
]
y = Inches(1.4)
for num, title, sub in items:
    add_rect(s, Inches(0.6), y, Inches(0.6), Inches(0.6), FABRIC_GREEN)
    add_text(s, Inches(0.6), y, Inches(0.6), Inches(0.6), num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.4), y - Inches(0.02), Inches(11.0), Inches(0.35), title, size=15, bold=True, color=INK)
    add_text(s, Inches(1.4), y + Inches(0.30), Inches(11.0), Inches(0.30), sub, size=11, color=GREY_MID)
    y += Inches(0.72)

# ============== SLIDE 3: CURRENT ARCHITECTURE (annotated) ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Manulife Japan today — current data architecture", "From the MLJ-Data-Architecture diagram", footer_idx=3)

# Top: sources row
row_x = Inches(0.5); row_w = SW - Inches(1.0)
y = Inches(1.3)
def layer_row(slide, label, y, items, fill=FABRIC_GREEN_LIGHT, border=FABRIC_GREEN, lbl_size=11):
    add_rect(slide, Inches(0.5), y, Inches(1.0), Inches(0.45), border)
    add_text(slide, Inches(0.5), y, Inches(1.0), Inches(0.45), label, size=lbl_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    n = len(items)
    avail = SW - Inches(1.7)
    cell_w = Emu(int(avail / n))
    x = Inches(1.6)
    for it in items:
        add_rect(slide, x, y, cell_w - Inches(0.05), Inches(0.45), fill, line=border, line_width=Pt(0.75))
        add_text(slide, x, y, cell_w - Inches(0.05), Inches(0.45), it, size=9, bold=False, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += cell_w

layer_row(s, "Source", Inches(1.25), ["Illustration & Application", "Underwriting", "Policy Hosting", "Policy Servicing", "Customer System", "Distributor"])
layer_row(s, "EDL", Inches(1.78), ["Data Mirror — Publish Database (single, shared mirror)"], fill=GREY_LIGHT, border=GREY_MID)
layer_row(s, "Cosmos", Inches(2.31), ["Document Data (Digital/AI)", "Customer Domain", "Distributor Domain", "Product Domain", "Finance Domain", "System Domain"])
layer_row(s, "Tables", Inches(2.84), ["Tabular Data (Analytics/AI)", "Customer Domain", "Distributor Domain", "Product Domain", "Finance Domain", "System Domain"], fill=GREY_LIGHT, border=GREY_MID)
layer_row(s, "Curated", Inches(3.37), ["Griffin", "Gobblin", "AML", "Purpose Driven Data", "IFRS", "CDP", "VOICE / CAR"])
layer_row(s, "Catalog", Inches(3.90), ["Unity Catalog  (governance — planned refresh 2026)"], fill=GREY_LIGHT, border=GREY_MID)
layer_row(s, "BI", Inches(4.43), ["Power BI Dashboards", "AI-Driven Data Self Service (planned 2026)"])
layer_row(s, "Users", Inches(4.96), ["Distribution", "IO Users", "Marketing", "Finance", "IT Users"], fill=GREY_LIGHT, border=GREY_MID)

# Observation box
add_round(s, Inches(0.5), Inches(5.7), SW - Inches(1.0), Inches(1.45), GREY_BG, line=GREY_LIGHT, line_width=Pt(0.5))
add_text(s, Inches(0.7), Inches(5.78), Inches(12.0), Inches(0.35), "What we observe", size=13, bold=True, color=FABRIC_GREEN)
add_bullets(s, Inches(0.7), Inches(6.10), SW - Inches(1.4), Inches(1.0), [
    "Clean domain-driven model — Customer / Distributor / Product / Finance / System — with separate document (Cosmos) and tabular paths.",
    "Multi-stack execution today: ADF for orchestration, Databricks for compute, Cosmos DB for documents, Unity Catalog for governance, Power BI for BI.",
    "Two strategic gaps already flagged for 2026: governance refresh on top of Unity Catalog, and an AI-driven self-service experience.",
], size=11, color=INK)

# ============== SLIDE 4: WHY FABRIC ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Why Fabric for Manulife Japan", "The technical case in one slide", footer_idx=4)
cards = [
    ("One copy of data",
     "OneLake replaces the EDL + Cosmos + flattened-tables shuffle with a single Delta lake. Every workload reads the same bytes — no Gobblin-style copy jobs to keep five layers in sync."),
    ("Domain ownership preserved",
     "Workspaces + lakehouses + Delta schemas map cleanly onto Customer / Distributor / Product / Finance / System. Each domain team owns its data; the platform stays unified."),
    ("Governance refresh, solved",
     "OneLake security + Purview unification gives you the 2026 governance step you've already planned — without porting Unity Catalog policies into a third tool."),
    ("AI-driven self-service, native",
     "Data Agent (NL → DAX/SQL/RAG) and Copilot for Power BI deliver the self-service AI box on your roadmap — wired to the governed semantic model out of the box."),
    ("Lower run cost & ops surface",
     "One capacity covers ingest, transform, warehouse, BI, and AI. Retire ADF, Databricks, and per-workspace SKU planning; pay for one F-SKU you can pause for non-business hours."),
    ("Open & non-locking",
     "Delta + Parquet on ADLS Gen2 underneath. Mirroring + Shortcuts let Cosmos DB, SQL DBs, and even Databricks tables show up in Fabric without movement."),
]
gx = Inches(0.5); gy = Inches(1.3); gap = Inches(0.2)
card_w = (SW - Inches(1.0) - gap*2) / 3
card_h = Inches(2.7)
for i, (h, body) in enumerate(cards):
    col = i % 3; row = i // 3
    x = gx + col * (card_w + gap)
    y = gy + row * (card_h + Inches(0.2))
    add_round(s, x, y, card_w, card_h, WHITE, line=FABRIC_GREEN, line_width=Pt(1.0))
    add_rect(s, x, y, Inches(0.15), card_h, FABRIC_GREEN)
    add_text(s, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.35), Inches(0.5),
             h, size=14, bold=True, color=FABRIC_GREEN)
    add_text(s, x + Inches(0.25), y + Inches(0.65), card_w - Inches(0.35), card_h - Inches(0.75),
             body, size=11, color=INK)

# ============== SLIDE 5: SIDE-BY-SIDE TABLE ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Side-by-side — current stack vs Fabric", "Layer-by-layer mapping", footer_idx=5)

# Table header
col_x = [Inches(0.5), Inches(2.5), Inches(7.0), Inches(7.0) + Inches(2.92)]
col_w = [Inches(2.0), Inches(4.5), Inches(2.92), Inches(2.93)]
table_top = Inches(1.3)
row_h = Inches(0.50)

# header row
headers = ["Layer", "Current (MLJ today)", "Fabric equivalent", "Why it's better"]
for i, htxt in enumerate(headers):
    fill = FABRIC_GREEN
    add_rect(s, col_x[i], table_top, col_w[i], row_h, fill)
    add_text(s, col_x[i] + Inches(0.1), table_top, col_w[i] - Inches(0.15), row_h,
             htxt, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("Source ingest",      "ADF + Databricks notebooks",                 "Data pipelines, Mirroring, Shortcuts",     "Zero-code Mirror for Cosmos/SQL; one tool"),
    ("EDL",                "Data Mirror Publish DB (relational copy)",   "OneLake bronze (Delta, schema-on-read)",   "No publish-DB serialisation; open format"),
    ("Doc data (Cosmos)",  "Cosmos DB per domain",                       "Mirrored Cosmos in OneLake + AI Search",   "Same Cosmos, queryable as Delta + indexed"),
    ("Tabular per domain", "Flattened tables in Databricks SQL",         "Lakehouse schemas per domain",             "Direct Lake — no import, no refresh"),
    ("Curated marts",      "Griffin / Gobblin / AML / IFRS / CDP / VOICE","Warehouse + semantic models on OneLake",  "Same marts, governed measures reused everywhere"),
    ("Governance",         "Unity Catalog (planned refresh 2026)",       "OneLake security + Microsoft Purview",     "The 2026 refresh, already integrated"),
    ("BI",                 "Power BI on Databricks SQL endpoint",        "Power BI on Direct Lake semantic model",   "Sub-second queries, no import refresh"),
    ("AI self-service",    "Planned for 2026",                           "Fabric Data Agent + Copilot",              "Delivers the box you've already drawn"),
    ("Compute SKUing",     "DBU $/hr + ADF IR + Cosmos RU + PBI Pro",    "Single F-SKU covers all workloads",        "One bill, pausable, smoothed bursts"),
]
y = table_top + row_h
for ridx, r in enumerate(rows):
    row_fill = WHITE if ridx % 2 == 0 else GREY_BG
    for i, val in enumerate(r):
        add_rect(s, col_x[i], y, col_w[i], row_h, row_fill, line=GREY_LIGHT, line_width=Pt(0.4))
        # accent border on Fabric column
        if i == 2:
            add_rect(s, col_x[i], y, Inches(0.06), row_h, FABRIC_GREEN)
        add_text(s, col_x[i] + Inches(0.12), y, col_w[i] - Inches(0.2), row_h,
                 val, size=10, color=INK if i != 1 else GREY_MID, bold=(i == 0), anchor=MSO_ANCHOR.MIDDLE)
    y += row_h

# ============== SLIDE 6: ARCHITECTURE — CURRENT (diagram) ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Current architecture — ADF + Databricks + Cosmos + Unity Catalog", subtitle="As reflected in MLJ-Data-Architecture.pdf", footer_idx=6)

def box(slide, x, y, w, h, text, fill=GREY_LIGHT, border=GREY_MID, color=INK, size=10, bold=False):
    add_rect(slide, x, y, w, h, fill, line=border, line_width=Pt(0.6))
    add_text(slide, x, y, w, h, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Sources
y0 = Inches(1.3)
box(s, Inches(0.5), y0, SW - Inches(1.0), Inches(0.45), "Source systems (6) — Illustration & App, Underwriting, Policy Hosting, Policy Servicing, Customer, Distributor", DATABRICKS_RED_LIGHT, DATABRICKS_RED, bold=True, size=11)

# ADF arrow
box(s, Inches(0.5), Inches(1.85), SW - Inches(1.0), Inches(0.4), "Azure Data Factory  —  orchestration + ingest", DATABRICKS_RED_LIGHT, DATABRICKS_RED, size=11)

# EDL Publish DB
box(s, Inches(0.5), Inches(2.35), SW - Inches(1.0), Inches(0.42), "EDL  —  Data Mirror Publish Database (relational copy of sources)", GREY_LIGHT, GREY_MID, size=11)

# Cosmos + Databricks parallel
box(s, Inches(0.5), Inches(2.90), Inches(6.0), Inches(0.45), "Cosmos DB  —  document data per domain (Customer / Distributor / Product / Finance / System)", DATABRICKS_RED_LIGHT, DATABRICKS_RED, size=10)
box(s, Inches(6.83), Inches(2.90), Inches(6.0), Inches(0.45), "Databricks  —  flattened tables per domain (analytics/AI)", DATABRICKS_RED_LIGHT, DATABRICKS_RED, size=10)

# Curated
box(s, Inches(0.5), Inches(3.5), SW - Inches(1.0), Inches(0.45), "Curated marts  —  Griffin · Gobblin · AML · Purpose-Driven · IFRS · CDP · VOICE/CAR", GREY_LIGHT, GREY_MID, size=10, bold=True)

# UC
box(s, Inches(0.5), Inches(4.05), Inches(8.0), Inches(0.45), "Unity Catalog  —  governance (refresh planned 2026)", DATABRICKS_RED_LIGHT, DATABRICKS_RED, size=10, bold=True)
box(s, Inches(8.83), Inches(4.05), Inches(4.0), Inches(0.45), "Power BI  —  dashboards", DATABRICKS_RED_LIGHT, DATABRICKS_RED, size=10, bold=True)

# Pain points box
add_round(s, Inches(0.5), Inches(4.85), SW - Inches(1.0), Inches(2.2), GREY_BG, line=GREY_LIGHT, line_width=Pt(0.4))
add_text(s, Inches(0.7), Inches(4.95), Inches(12), Inches(0.4), "Friction points in the current stack", size=13, bold=True, color=DATABRICKS_RED)
add_bullets(s, Inches(0.7), Inches(5.35), SW - Inches(1.4), Inches(1.7), [
    "Three platforms to govern (Databricks UC, Cosmos RBAC, Power BI workspace ACLs) — three audit surfaces, three SKUs.",
    "Document layer (Cosmos) is separated from tabular layer — RAG/AI use cases need bespoke glue code per domain.",
    "Multiple physical copies: Source → Publish DB → Cosmos / Tables → Curated. Every hop is latency, cost, and a divergence risk.",
    "Power BI imports from Databricks SQL — refresh windows, dataset size limits, and stale data in dashboards.",
    "AI-driven self-service still on the 2026 roadmap; no native NL-to-SQL grounded on the semantic layer today.",
], size=11, color=INK)

# ============== SLIDE 7: ARCHITECTURE — FABRIC TARGET ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Target architecture — Microsoft Fabric for MLJ", subtitle="Same domain model, one platform, governance-first", footer_idx=7)

# Sources
y0 = Inches(1.3)
box(s, Inches(0.5), y0, SW - Inches(1.0), Inches(0.42), "Source systems (6) — unchanged", FABRIC_GREEN_LIGHT, FABRIC_GREEN, size=11, bold=True)

# Ingest
box(s, Inches(0.5), Inches(1.78), SW - Inches(1.0), Inches(0.42), "Fabric Data Pipelines  +  Mirroring  +  Shortcuts   (replace ADF, no separate orchestration tier)", FABRIC_GREEN_LIGHT, FABRIC_GREEN, size=11)

# OneLake bronze
box(s, Inches(0.5), Inches(2.28), SW - Inches(1.0), Inches(0.42), "OneLake — Bronze (Delta, schema-on-read)   ← Mirrors source DBs in near real-time", GREY_LIGHT, GREY_MID, size=11)

# Silver per-domain
silver_x = Inches(0.5)
n=5; total_w = SW - Inches(1.0); cw = Emu(int(total_w / n))
domains = ["Customer", "Distributor", "Product", "Finance", "System"]
for i, d in enumerate(domains):
    box(s, silver_x + i*cw, Inches(2.78), cw - Inches(0.05), Inches(0.42), f"Silver — {d}", FABRIC_GREEN_LIGHT, FABRIC_GREEN, size=10)

# Gold marts row
box(s, Inches(0.5), Inches(3.28), SW - Inches(1.0), Inches(0.45),
    "Gold — Curated marts (Griffin · AML · IFRS · CDP · VOICE/CAR) as Warehouses / semantic models on OneLake",
    GREY_LIGHT, GREY_MID, size=10, bold=True)

# Doc layer integrated
box(s, Inches(0.5), Inches(3.83), Inches(6.0), Inches(0.45), "Document data — same Cosmos DB, mirrored into OneLake + indexed in Azure AI Search", FABRIC_GREEN_LIGHT, FABRIC_GREEN, size=10)
box(s, Inches(6.83), Inches(3.83), Inches(6.0), Inches(0.45), "Direct Lake Semantic Model — single source of measures", FABRIC_GREEN_LIGHT, FABRIC_GREEN, size=10, bold=True)

# Governance + Consumption
box(s, Inches(0.5), Inches(4.4), Inches(7.0), Inches(0.45), "Governance — OneLake security + Microsoft Purview (replaces the planned Unity Catalog refresh)", FABRIC_GREEN_LIGHT, FABRIC_GREEN, size=10, bold=True)
box(s, Inches(7.83), Inches(4.4), Inches(5.0), Inches(0.45), "Consumption — Power BI · Data Agent · Copilot · APIs", FABRIC_GREEN_LIGHT, FABRIC_GREEN, size=10, bold=True)

# Benefits
add_round(s, Inches(0.5), Inches(5.2), SW - Inches(1.0), Inches(1.85), GREY_BG, line=GREY_LIGHT, line_width=Pt(0.4))
add_text(s, Inches(0.7), Inches(5.28), Inches(12), Inches(0.4), "What changes", size=13, bold=True, color=FABRIC_GREEN)
add_bullets(s, Inches(0.7), Inches(5.62), SW - Inches(1.4), Inches(1.45), [
    "One platform (Fabric) replaces ADF + Databricks + the bespoke serving + Power BI Premium — one bill, pausable.",
    "Documents (Cosmos) live next to tabular data in OneLake — RAG + analytics share the same governance and lineage.",
    "Direct Lake means Power BI and Data Agent query Delta in place — no dataset import, no refresh window, no stale dashboards.",
    "The 2026 governance refresh box is filled by OneLake + Purview without a third tool migration.",
], size=11, color=INK)

# ============== SLIDE 8: ONELAKE FOUNDATION ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "OneLake — the unified data foundation", "One copy, open format, shared across every Fabric workload", footer_idx=8)

# Left: diagram
add_round(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.4), GREY_BG, line=GREY_LIGHT, line_width=Pt(0.5))
add_text(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(0.4), "How OneLake replaces multi-store sprawl", size=13, bold=True, color=FABRIC_GREEN)

# big OneLake circle
add_rect(s, Inches(2.4), Inches(2.5), Inches(2.4), Inches(2.4), FABRIC_GREEN)
add_text(s, Inches(2.4), Inches(2.5), Inches(2.4), Inches(2.4), "OneLake\n(Delta / Parquet on ADLS Gen2)", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# surrounding consumers
consumers = [
    ("Spark notebooks", Inches(0.7), Inches(2.1)),
    ("Warehouse (T-SQL)", Inches(5.0), Inches(2.1)),
    ("Power BI", Inches(0.7), Inches(3.4)),
    ("Data Agent / AI", Inches(5.0), Inches(3.4)),
    ("KQL Eventhouse", Inches(0.7), Inches(4.7)),
    ("Mirrored Cosmos / SQL", Inches(5.0), Inches(4.7)),
]
for name, x, y in consumers:
    add_round(s, x, y, Inches(1.5), Inches(0.5), WHITE, line=FABRIC_GREEN, line_width=Pt(0.8))
    add_text(s, x, y, Inches(1.5), Inches(0.5), name, size=10, bold=True, color=FABRIC_GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.7), Inches(5.4), Inches(6.0), Inches(1.3),
         "All consumers read the same Delta tables — no copies, no exports, no ETL-to-BI step. Mirroring keeps Cosmos/SQL data in sync without movement code.",
         size=10, color=GREY_MID)

# Right: benefits for MLJ
add_round(s, Inches(7.0), Inches(1.4), Inches(5.83), Inches(5.4), WHITE, line=FABRIC_GREEN, line_width=Pt(1.0))
add_rect(s, Inches(7.0), Inches(1.4), Inches(0.15), Inches(5.4), FABRIC_GREEN)
add_text(s, Inches(7.25), Inches(1.5), Inches(5.5), Inches(0.4), "Why this matters for MLJ", size=13, bold=True, color=FABRIC_GREEN)
add_bullets(s, Inches(7.25), Inches(1.95), Inches(5.5), Inches(4.7), [
    "One copy of data — retire the Source → Publish DB → Cosmos/Tables chain. Same domains, one physical layer.",
    "Open format — Delta + Parquet over ADLS Gen2. Non-Microsoft tools can still read it.",
    "Mirroring — Cosmos DB and SQL stay where they are. OneLake gets a near-real-time Delta view.",
    "Shortcuts — surface S3, ADLS, or Databricks-managed tables in OneLake without copy.",
    "Workspace-aligned to domains — Customer, Distributor, Product, Finance, System.",
    "Sensitivity labels & lineage propagate end-to-end — no per-tool re-labelling.",
    "Pausable F-SKU billing — no storage cost on top, storage is paid as ADLS Gen2.",
], size=11, color=INK)

# ============== SLIDE 9: UNITY CATALOG vs FABRIC GOVERNANCE ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Unity Catalog → OneLake + Purview", "Filling the governance-refresh box you've already planned for 2026", footer_idx=9)

# Two columns: UC | Fabric
def panel(slide, x, y, w, h, title, color):
    add_rect(slide, x, y, w, Inches(0.5), color)
    add_text(slide, x, y, w, Inches(0.5), title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, x, y + Inches(0.5), w, h - Inches(0.5), WHITE, line=color, line_width=Pt(1.0))

# Left UC
panel(s, Inches(0.5), Inches(1.3), Inches(6.2), Inches(5.5), "Unity Catalog (today)", DATABRICKS_RED)
add_bullets(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.7), [
    "Governance scope: Databricks workspaces + Delta tables registered to UC.",
    "Cosmos DB, Power BI, ADF, source systems live OUTSIDE UC — separate policies.",
    "Lineage: strong inside Databricks; gaps at boundaries (Cosmos in, Power BI out).",
    "Sensitivity labels: not native — needs Purview side-by-side anyway.",
    "Audit surface: UC + Cosmos RBAC + PBI workspace ACLs + Azure RBAC = 4 places to look.",
    "Refresh planned for 2026 — that work is already on the roadmap.",
], size=11, color=INK)

# Right Fabric
panel(s, Inches(6.83), Inches(1.3), Inches(6.0), Inches(5.5), "OneLake + Microsoft Purview (target)", FABRIC_GREEN)
add_bullets(s, Inches(7.03), Inches(1.9), Inches(5.7), Inches(4.7), [
    "Governance scope: everything in OneLake + Mirrored sources + Power BI items + Data Agents — one catalog.",
    "Built-in sensitivity labels, DLP, and access policies that propagate to consumers automatically.",
    "End-to-end lineage from source mirror through bronze/silver/gold to PBI report and Data Agent answer.",
    "Purview unifies discovery across Fabric, Azure data services, and on-prem — one audit surface.",
    "OneLake security: row, column, and table policies enforced wherever data is read.",
    "Replaces the planned UC refresh — and delivers AI-readiness for grounding & responsible AI controls.",
], size=11, color=INK)

# ============== SLIDE 10: AI-DRIVEN SELF-SERVICE ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "AI-driven data self-service", "Fabric Data Agent + Copilot — your 2026 AI box, available today", footer_idx=10)

# Conceptual flow
add_round(s, Inches(0.5), Inches(1.4), SW - Inches(1.0), Inches(2.0), GREY_BG, line=GREY_LIGHT, line_width=Pt(0.5))
flow = ["Business user asks in JP/EN", "Data Agent (grounded)", "Semantic model + Lakehouse + AI Search", "Answer with citations"]
n=len(flow); fw=Inches(2.8); gap=Inches(0.2); start_x = Inches(0.9)
y = Inches(2.05); h = Inches(0.7)
for i, t in enumerate(flow):
    x = start_x + i * (fw + gap)
    add_round(s, x, y, fw, h, FABRIC_GREEN, line=None)
    add_text(s, x, y, fw, h, t, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < n-1:
        ax = x + fw; ay = y + h/2
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay - Inches(0.1), gap, Inches(0.2))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = GREY_MID
        arrow.line.fill.background()

add_text(s, Inches(0.7), Inches(2.95), SW - Inches(1.4), Inches(0.4),
         "Grounding is constrained to the curated semantic model + selected lakehouse tables + AI Search index — answers cite the source.",
         size=10, color=GREY_MID, align=PP_ALIGN.CENTER)

# Use cases panel
add_round(s, Inches(0.5), Inches(3.6), Inches(6.2), Inches(3.4), WHITE, line=FABRIC_GREEN, line_width=Pt(1.0))
add_text(s, Inches(0.7), Inches(3.7), Inches(6.0), Inches(0.4), "Use cases at MLJ", size=13, bold=True, color=FABRIC_GREEN)
add_bullets(s, Inches(0.7), Inches(4.15), Inches(6.0), Inches(2.8), [
    "Distribution: \"Top 10 advisors by AUM in Tokyo this quarter\"",
    "Underwriting: \"Show declined applications by reason, last 90 days\"",
    "Finance: \"IFRS 17 contractual service margin run-off by product\"",
    "Marketing: \"Variable annuity policyholders 55+ with no recent contact\"",
    "Customer service: \"What does the policy say about disability rider exclusions?\"",
    "Compliance: \"AML alerts open more than 30 days by branch\"",
], size=11, color=INK)

# Differentiators
add_round(s, Inches(7.0), Inches(3.6), Inches(5.83), Inches(3.4), WHITE, line=FABRIC_GREEN, line_width=Pt(1.0))
add_text(s, Inches(7.2), Inches(3.7), Inches(5.6), Inches(0.4), "What makes it work for MLJ", size=13, bold=True, color=FABRIC_GREEN)
add_bullets(s, Inches(7.2), Inches(4.15), Inches(5.6), Inches(2.8), [
    "Governed grounding — answers come from curated semantic model, not free text.",
    "Citations — every answer includes the table or document source.",
    "Japanese-language ready — agent instructions and embeddings support JP.",
    "Inherits OneLake security — users see only what they're entitled to.",
    "No new tool to deploy — runs inside Fabric, same workspace as the data.",
], size=11, color=INK)

# ============== SLIDE 11: DOMAIN MAPPING ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Domain mapping — MLJ on Fabric", "How the 5 MLJ domains land as workspaces, lakehouses, and marts", footer_idx=11)

domains_full = [
    ("Customer",     "Customer System,\nIllustration & App",      "Policies, applications, contacts",        "CDP, VOICE/CAR"),
    ("Distributor",  "Distributor",                                  "Advisors, agencies, performance",         "Distributor 360"),
    ("Product",      "Policy Hosting,\nProduct Master",             "Product catalogue, riders, illustrations","Product P&L"),
    ("Finance",      "Policy Servicing,\nGL feeds",                  "Premiums, claims, GL postings",           "IFRS 17, CSM, P&L"),
    ("System",       "Application logs,\nops telemetry",             "DQ flags, batch status, lineage",         "Griffin, AML"),
]
# header row
hx = Inches(0.5); hy = Inches(1.3)
cols = ["Domain", "Source systems", "Silver tables (per domain)", "Gold curated marts"]
cw = [Inches(1.6), Inches(3.0), Inches(4.3), Inches(3.43)]
for i, c in enumerate(cols):
    x = hx + sum(cw[:i])
    add_rect(s, x, hy, cw[i], Inches(0.5), FABRIC_GREEN)
    add_text(s, x + Inches(0.1), hy, cw[i] - Inches(0.15), Inches(0.5), c, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
y = hy + Inches(0.5)
for ridx, dom in enumerate(domains_full):
    row_fill = WHITE if ridx % 2 == 0 else GREY_BG
    for i, val in enumerate(dom):
        x = hx + sum(cw[:i])
        add_rect(s, x, y, cw[i], Inches(0.85), row_fill, line=GREY_LIGHT, line_width=Pt(0.4))
        add_text(s, x + Inches(0.12), y, cw[i] - Inches(0.2), Inches(0.85), val,
                 size=11, color=INK, bold=(i == 0), anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.85)

add_round(s, Inches(0.5), Inches(6.4), SW - Inches(1.0), Inches(0.7), GREY_BG, line=GREY_LIGHT, line_width=Pt(0.4))
add_text(s, Inches(0.7), Inches(6.45), SW - Inches(1.4), Inches(0.6),
         "Each MLJ domain → a Fabric workspace (or schema) with its own lakehouse, owned by the domain team. Platform governance lives above, via OneLake + Purview.",
         size=11, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# ============== SLIDE 12: DEMO FLOW ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Live demo — what we'll show", "Built in the ManulifeJapan-Fabric-POC workspace on F64 capacity", footer_idx=12)

steps = [
    ("1. OneLake & domain layout",
     "Show the 5 domain lakehouses with Bronze/Silver/Gold Delta tables. Cosmos-style document table side-by-side with tabular."),
    ("2. Notebook orchestration",
     "Run end-to-end Spark notebooks: ingest 6 sources → 5 domain silvers → curated marts (Griffin / AML / IFRS). One platform, no ADF."),
    ("3. Direct Lake semantic model",
     "Queries hit Delta in place — no import, no refresh window. JPY measures: Total Premium, AUM, Claims Ratio, Approval Rate."),
    ("4. Power BI dashboard",
     "A Power BI report against the semantic model showing premium by prefecture and product, fed live by Direct Lake."),
    ("5. Fabric Data Agent",
     "Ask in English & Japanese: \"Top 5 advisors by AUM in Tokyo\" / \"variable annuity claims approval rate this quarter\". Cites tables + documents."),
    ("6. Governance & lineage",
     "OneLake security & Purview lineage — single audit surface across ingest → mart → Power BI → Data Agent."),
]
y = Inches(1.3)
for title, body in steps:
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.85), FABRIC_GREEN)
    add_text(s, Inches(0.75), y + Inches(0.05), Inches(12.0), Inches(0.35), title, size=14, bold=True, color=FABRIC_GREEN)
    add_text(s, Inches(0.75), y + Inches(0.40), Inches(12.0), Inches(0.45), body, size=11, color=INK)
    y += Inches(0.92)

# ============== SLIDE 13: MIGRATION PHASING ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Migration phasing — pragmatic path from current to target", footer_idx=13)

phases = [
    ("Phase 0\nDiscover",       "1-2 wks", "Inventory ADF pipelines, Databricks jobs, UC schemas, Cosmos collections.\nIdentify quick wins (Mirroring candidates).", FABRIC_GREEN_LIGHT),
    ("Phase 1\nCoexist",        "6-8 wks", "Stand up OneLake. Mirror Cosmos + SQL. Shortcut Databricks-managed tables. Build first domain lakehouse (Customer) + semantic model. Parallel run.", FABRIC_GREEN_LIGHT),
    ("Phase 2\nMigrate",        "3-6 mo",  "Migrate domain-by-domain: Customer → Distributor → Product → Finance → System. Re-platform curated marts (Griffin, AML, IFRS). Power BI moves to Direct Lake.", FABRIC_GREEN_LIGHT),
    ("Phase 3\nRetire",         "1-2 mo",  "Decommission ADF pipelines, Databricks jobs, UC. Move governance fully to OneLake + Purview. Roll out Data Agent + Copilot.", FABRIC_GREEN_LIGHT),
    ("Phase 4\nOptimise",       "ongoing", "Capacity right-sizing, autoscaling, pause windows, monthly FinOps review. Expand AI use cases.", FABRIC_GREEN_LIGHT),
]
gx = Inches(0.5); gy = Inches(1.4); gap = Inches(0.18)
n = len(phases)
total_w = SW - Inches(1.0) - gap*(n-1)
phase_w = Emu(int(total_w / n))
phase_h = Inches(4.5)
for i, (ttl, dur, body, fill) in enumerate(phases):
    x = gx + i * (phase_w + gap)
    add_round(s, x, gy, phase_w, phase_h, WHITE, line=FABRIC_GREEN, line_width=Pt(1.0))
    add_rect(s, x, gy, phase_w, Inches(0.9), FABRIC_GREEN)
    add_text(s, x, gy + Inches(0.05), phase_w, Inches(0.6), ttl, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, gy + Inches(0.55), phase_w, Inches(0.35), dur, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), gy + Inches(1.0), phase_w - Inches(0.2), phase_h - Inches(1.1), body, size=10, color=INK)

add_round(s, Inches(0.5), Inches(6.15), SW - Inches(1.0), Inches(0.9), GREY_BG, line=GREY_LIGHT, line_width=Pt(0.4))
add_text(s, Inches(0.7), Inches(6.20), SW - Inches(1.4), Inches(0.4), "Key principle — never break production", size=12, bold=True, color=FABRIC_GREEN)
add_text(s, Inches(0.7), Inches(6.50), SW - Inches(1.4), Inches(0.5),
         "All migration happens in parallel run mode — Fabric outputs are validated against the current stack before any cutover. The 2026 governance refresh moves under one umbrella instead of being a separate workstream.",
         size=10, color=INK)

# ============== SLIDE 14: NEXT STEPS ==============
s = prs.slides.add_slide(BLANK)
slide_header(s, "Next steps", footer_idx=14)

cards = [
    ("This week", "Walk through the POC together; agree which MLJ domain runs Phase 1 (recommend Customer).", FABRIC_GREEN),
    ("Within 30 days", "Mirroring proof: stand up a real Cosmos DB mirror + a Direct Lake report on one production-shape table.", FABRIC_GREEN),
    ("Within 60 days", "Phase 1 design: workspace topology, capacity sizing, security model, Purview rollout.", FABRIC_GREEN),
    ("Within 90 days", "Phase 1 build kick-off; parallel-run plan for Customer domain.", FABRIC_GREEN),
]
y = Inches(1.5)
for i, (k, v, c) in enumerate(cards):
    add_round(s, Inches(0.7), y, Inches(2.3), Inches(1.0), c)
    add_text(s, Inches(0.7), y, Inches(2.3), Inches(1.0), k, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_round(s, Inches(3.2), y, Inches(9.6), Inches(1.0), WHITE, line=c, line_width=Pt(1.0))
    add_text(s, Inches(3.35), y, Inches(9.5), Inches(1.0), v, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.15)

add_text(s, Inches(0.5), Inches(6.5), SW - Inches(1.0), Inches(0.5),
         "Thank you — questions?", size=22, bold=True, color=FABRIC_GREEN, align=PP_ALIGN.CENTER)

# ------------------ save ------------------
out_path = r"C:\Users\anuragdhuria\Downloads\Manulife Japan\Manulife-Japan-Fabric-POC.pptx"
prs.save(out_path)
print(f"Saved {out_path}")
